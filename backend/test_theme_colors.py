"""
Test rápido para verificar extracción de colores del tema
"""
from pptx import Presentation
from pptx_analyzer import get_theme_colors, extract_background
import sys

def test_theme_extraction(pptx_path):
    """Prueba la extracción de colores del tema"""
    print(f"\n{'='*80}")
    print(f"TEST: Extracción de Colores del Tema")
    print(f"{'='*80}\n")
    
    try:
        prs = Presentation(pptx_path)
        
        print(f"📄 Archivo: {pptx_path}")
        print(f"📊 Total de slides: {len(prs.slides)}\n")
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n--- SLIDE {i} ---")
            
            # Extraer colores del tema
            theme_colors = get_theme_colors(slide)
            print(f"🎨 Colores del tema extraídos: {len(theme_colors)}")
            for name, color in theme_colors.items():
                print(f"   {name}: {color}")
            
            # Extraer fondo
            print(f"\n🖼️ Extrayendo fondo...")
            background = extract_background(slide)
            print(f"   Tipo: {background['type']}")
            print(f"   Color: {background['color']}")
            
        print(f"\n{'='*80}")
        print("✅ Test completado")
        print(f"{'='*80}\n")
        
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    # Buscar archivo PPTX en el directorio actual o temp
    import os
    import glob
    
    # Buscar en temp
    temp_files = glob.glob(r"C:\Users\admin\AppData\Local\Temp\tmp*.pptx")
    
    if temp_files:
        # Usar el más reciente
        latest = max(temp_files, key=os.path.getmtime)
        print(f"📁 Usando archivo temporal: {latest}")
        test_theme_extraction(latest)
    else:
        print("❌ No se encontró archivo PPTX en temp")
        print("💡 Sube un template en la app primero")

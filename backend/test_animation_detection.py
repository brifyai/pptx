"""
Test script para verificar la detección de animaciones
"""
import sys
from pptx import Presentation
from pptx_analyzer import detect_animated_shapes, extract_all_assets

def test_animation_detection(pptx_path):
    """Prueba la detección de animaciones en un PPTX"""
    print(f"\n{'='*60}")
    print(f"PROBANDO DETECCIÓN DE ANIMACIONES")
    print(f"Archivo: {pptx_path}")
    print(f"{'='*60}\n")
    
    try:
        prs = Presentation(pptx_path)
        
        print(f"📄 Slides encontrados: {len(prs.slides)}\n")
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\n{'─'*60}")
            print(f"SLIDE {slide_idx + 1}")
            print(f"{'─'*60}")
            
            # Detectar animaciones
            animated_ids = detect_animated_shapes(slide)
            
            if animated_ids:
                print(f"\n✅ Animaciones detectadas: {len(animated_ids)} shapes")
                for shape_id in animated_ids:
                    print(f"   - Shape ID: {shape_id}")
            else:
                print(f"\n❌ No se detectaron animaciones")
            
            # Mostrar información de shapes
            print(f"\n📦 Shapes en el slide:")
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    is_animated = shape.shape_id in animated_ids
                    print(f"   - Shape ID {shape.shape_id}: PICTURE {'🎬 ANIMADO' if is_animated else ''}")
                    
                    # Info adicional
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    width_ratio = shape.width / slide_width
                    height_ratio = shape.height / slide_height
                    print(f"     Tamaño: {width_ratio:.1%} x {height_ratio:.1%} del slide")
        
        # Probar extract_all_assets
        print(f"\n{'='*60}")
        print(f"PROBANDO EXTRACCIÓN DE ASSETS")
        print(f"{'='*60}\n")
        
        assets = extract_all_assets(prs)
        
        print(f"\n📊 RESUMEN:")
        print(f"   Total: {assets['totalCount']}")
        print(f"   Logos: {len(assets['logos'])}")
        print(f"   Transparentes: {len(assets['transparentImages'])}")
        print(f"   Animados: {len(assets['animatedElements'])}")
        print(f"   Imágenes: {len(assets['images'])}")
        
        if assets['animatedElements']:
            print(f"\n🎬 ELEMENTOS ANIMADOS:")
            for asset in assets['animatedElements']:
                print(f"   - Slide {asset['slideNumber']}, Shape {asset['shapeId']}")
                print(f"     Logo: {asset['isLogo']}, Transparencia: {asset['hasTransparency']}")
        
    except Exception as e:
        print(f"\n❌ ERROR: {e}")
        import traceback
        traceback.print_exc()

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Uso: python test_animation_detection.py <template.pptx>")
        print("\nEjemplo:")
        print("  python test_animation_detection.py uploads/template.pptx")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    test_animation_detection(pptx_path)

from fastapi import FastAPI, UploadFile, File, HTTPException, Body, WebSocket, WebSocketDisconnect, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
import uvicorn
from pptx_analyzer import analyze_presentation
from pptx_generator import generate_presentation
from database import PresentationDB
from mapping_cache import MappingCache
from shape_matcher import ShapeMatcher
from font_detector import analyze_fonts as analyze_pptx_fonts
from typing import Dict, Any, List, Optional
import os
import tempfile
import json
from pptx import Presentation
from pptx.util import Inches, Pt
import subprocess
import asyncio
import httpx
import base64
import logging
import traceback
from datetime import datetime

# Configure logging with detailed format (Requirement 6.5)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

# Error categories for structured logging (Requirement 6.5)
class ErrorCategory:
    FILE_VALIDATION = "FILE_VALIDATION"
    FILE_CONVERSION = "FILE_CONVERSION"
    GEMINI_API = "GEMINI_API"
    SHAPE_MATCHING = "SHAPE_MATCHING"
    CACHE_OPERATION = "CACHE_OPERATION"
    DATABASE = "DATABASE"
    NETWORK = "NETWORK"
    UNKNOWN = "UNKNOWN"


def log_error_with_context(
    error: Exception,
    category: str,
    operation: str,
    context: Dict[str, Any] = None,
    include_traceback: bool = True
):
    """
    Log an error with full context for diagnosis (Requirement 6.5).
    
    Args:
        error: The exception that occurred
        category: Error category from ErrorCategory
        operation: The operation that was being performed
        context: Additional context information
        include_traceback: Whether to include full traceback
    """
    error_info = {
        "timestamp": datetime.utcnow().isoformat(),
        "category": category,
        "operation": operation,
        "error_type": type(error).__name__,
        "error_message": str(error),
        "context": context or {}
    }
    
    # Log the structured error
    logger.error(f"❌ [{category}] {operation} failed: {str(error)}")
    logger.error(f"   Context: {json.dumps(error_info['context'], default=str)}")
    
    if include_traceback:
        tb = traceback.format_exc()
        logger.error(f"   Traceback:\n{tb}")
    
    return error_info


def log_operation_start(operation: str, context: Dict[str, Any] = None):
    """Log the start of an operation with context."""
    logger.info(f"🔄 Starting: {operation}")
    if context:
        logger.info(f"   Context: {json.dumps(context, default=str)}")


def log_operation_success(operation: str, result_summary: str = None):
    """Log successful completion of an operation."""
    msg = f"✅ Completed: {operation}"
    if result_summary:
        msg += f" - {result_summary}"
    logger.info(msg)

app = FastAPI(title="AI Presentation API")
db = PresentationDB()
mapping_cache = MappingCache()

# Aumentar límite de tamaño de archivo (100MB)
from starlette.middleware import Middleware
from starlette.requests import Request

# Configurar límite de tamaño máximo para uploads
MAX_UPLOAD_SIZE = 100 * 1024 * 1024  # 100MB

# Pydantic models for request/response validation
class UpdateMappingRequest(BaseModel):
    template_hash: str
    element_id: str
    new_type: str

class UpdateMappingResponse(BaseModel):
    success: bool
    message: str

class AnalysisElement(BaseModel):
    id: str
    type: str
    coordinates: Dict[str, int]
    style: Optional[Dict[str, Any]] = None
    confidence: Optional[float] = None
    shapeId: Optional[int] = None

class AnalysisResponse(BaseModel):
    success: bool
    templateHash: str
    elements: List[AnalysisElement]
    shapeMapping: Dict[str, int]
    source: str  # 'cache' or 'vision'
    message: str

# Gestión de conexiones WebSocket
class ConnectionManager:
    def __init__(self):
        self.active_connections: Dict[str, List[WebSocket]] = {}
    
    async def connect(self, websocket: WebSocket, presentation_id: str):
        await websocket.accept()
        if presentation_id not in self.active_connections:
            self.active_connections[presentation_id] = []
        self.active_connections[presentation_id].append(websocket)
        print(f"✅ Cliente conectado a presentación {presentation_id}")
    
    def disconnect(self, websocket: WebSocket, presentation_id: str):
        if presentation_id in self.active_connections:
            self.active_connections[presentation_id].remove(websocket)
            if not self.active_connections[presentation_id]:
                del self.active_connections[presentation_id]
        print(f"❌ Cliente desconectado de presentación {presentation_id}")
    
    async def broadcast(self, presentation_id: str, message: dict, exclude: WebSocket = None):
        if presentation_id in self.active_connections:
            disconnected = []
            for connection in self.active_connections[presentation_id]:
                if connection != exclude:
                    try:
                        await connection.send_json(message)
                    except:
                        disconnected.append(connection)
            
            # Limpiar conexiones muertas
            for conn in disconnected:
                self.disconnect(conn, presentation_id)

manager = ConnectionManager()

# Configurar límite de tamaño para uploads (50MB)
from starlette.middleware.base import BaseHTTPMiddleware

class LargeUploadMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request, call_next):
        # Permitir uploads grandes
        return await call_next(request)

# CORS para permitir requests desde React
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3006", "http://localhost:3007", "http://localhost:3008", "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
async def root():
    return {
        "message": "AI Presentation API",
        "version": "1.0.0",
        "endpoints": {
            "analyze": "POST /api/analyze - Analiza un PPT y extrae su diseño",
            "generate": "POST /api/generate - Genera PPT con contenido de IA"
        }
    }

@app.post("/api/analyze")
async def analyze_ppt(file: UploadFile = File(...)):
    """
    Analiza un archivo PowerPoint y extrae toda su estructura de diseño
    """
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    try:
        # Guardar archivo temporalmente
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name
        
        print(f"📄 Archivo guardado en: {tmp_path}")
        
        # Analizar presentación
        analysis = analyze_presentation(tmp_path)
        
        print(f"✅ Análisis completado: {len(analysis['slides'])} slides")
        
        # Limpiar archivo temporal
        os.unlink(tmp_path)
        
        return {
            "success": True,
            "analysis": analysis,
            "message": f"Análisis completado: {len(analysis['slides'])} diapositivas detectadas"
        }
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al analizar: {str(e)}")

@app.post("/api/generate")
async def generate_ppt(
    file: UploadFile = File(...),
    content: str = None
):
    """
    Genera un nuevo PowerPoint manteniendo el diseño original
    pero con contenido generado por IA
    """
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    try:
        # Guardar archivo original temporalmente
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            file_content = await file.read()
            tmp.write(file_content)
            original_path = tmp.name
        
        # Parsear contenido generado por IA
        if content:
            ai_content = json.loads(content)
        else:
            ai_content = None
        
        # Generar nueva presentación
        output_path = generate_presentation(original_path, ai_content)
        
        # Limpiar archivo original
        os.unlink(original_path)
        
        # Devolver archivo generado
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="presentacion_generada.pptx",
            headers={"Content-Disposition": "attachment; filename=presentacion_generada.pptx"}
        )
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al generar: {str(e)}")

@app.get("/health")
async def health_check():
    return {"status": "healthy", "service": "AI Presentation API"}

@app.post("/api/export/pptx")
async def export_pptx(request: Request):
    """
    Exporta slides a PowerPoint (.pptx)
    Si se proporciona un template, clona su diseño y solo reemplaza el texto.
    Si no hay template, crea una presentación básica.
    
    Acepta archivos grandes (hasta 50MB).
    """
    try:
        # Configurar límite de tamaño para este request (50MB)
        from starlette.datastructures import UploadFile as StarletteUploadFile
        
        # Leer el formulario con límite aumentado
        # python-multipart usa max_file_size por defecto de 1MB
        # Necesitamos leer manualmente
        
        content_type = request.headers.get('content-type', '')
        
        if 'multipart/form-data' in content_type:
            # Leer body completo
            body = await request.body()
            
            # Parsear manualmente el multipart
            import re
            from io import BytesIO
            
            # Extraer boundary
            boundary_match = re.search(r'boundary=([^;]+)', content_type)
            if not boundary_match:
                raise HTTPException(status_code=400, detail="Invalid multipart boundary")
            
            boundary = boundary_match.group(1).strip('"')
            
            # Parsear partes
            template_content = None
            template_filename = None
            data = None
            
            parts = body.split(f'--{boundary}'.encode())
            
            for part in parts:
                if b'name="template"' in part:
                    # Extraer filename
                    filename_match = re.search(rb'filename="([^"]+)"', part)
                    if filename_match:
                        template_filename = filename_match.group(1).decode('utf-8')
                    
                    # Extraer contenido (después de \r\n\r\n)
                    content_start = part.find(b'\r\n\r\n')
                    if content_start != -1:
                        template_content = part[content_start + 4:]
                        # Remover trailing \r\n
                        if template_content.endswith(b'\r\n'):
                            template_content = template_content[:-2]
                
                elif b'name="data"' in part:
                    content_start = part.find(b'\r\n\r\n')
                    if content_start != -1:
                        data_bytes = part[content_start + 4:]
                        if data_bytes.endswith(b'\r\n'):
                            data_bytes = data_bytes[:-2]
                        data = data_bytes.decode('utf-8')
            
            logger.info(f"📤 Export PPTX - Template: {template_filename}")
            logger.info(f"📤 Export PPTX - Template size: {len(template_content) if template_content else 0} bytes")
            logger.info(f"📤 Export PPTX - Data: {data[:200] if data else 'None'}...")
        else:
            # JSON request
            body = await request.body()
            json_data = json.loads(body)
            data = json.dumps(json_data)
            template_content = None
            template_filename = None
        
        # Parsear datos de slides
        slides = []
        if data:
            slides_data = json.loads(data) if isinstance(data, str) else data
            slides = slides_data.get('slides', [])
            logger.info(f"📤 Slides parseados: {len(slides)}")
        
        # Si hay template, usar clonación completa
        if template_content and template_filename:
            logger.info(f"📄 Usando template: {template_filename}")
            
            # Guardar template temporalmente
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                tmp.write(template_content)
                template_path = tmp.name
            
            logger.info(f"📄 Template guardado en: {template_path} ({len(template_content)} bytes)")
            
            # Preparar contenido en formato esperado por generate_presentation
            ai_content = {
                'slides': [{'content': slide.get('content', {})} for slide in slides]
            }
            
            logger.info(f"📝 Contenido preparado para {len(ai_content['slides'])} slides")
            for idx, slide in enumerate(ai_content['slides']):
                content = slide.get('content', {})
                logger.info(f"   Slide {idx+1}: {list(content.keys())}")
            
            # Generar usando clonación completa
            output_path = generate_presentation(template_path, ai_content)
            
            logger.info(f"✅ Presentación generada: {output_path}")
            
            # Limpiar template temporal
            os.unlink(template_path)
            
            return FileResponse(
                output_path,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                filename="presentacion.pptx",
                headers={"Content-Disposition": "attachment; filename=presentacion.pptx"}
            )
        
        # Si no hay template, crear presentación básica (fallback)
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
        
        for slide_data in slides:
            slide_type = slide_data.get('type', 'content')
            content = slide_data.get('content', {})
            
            if slide_type == 'title':
                # Layout de título
                layout = prs.slide_layouts[6]  # Blank
                slide = prs.slides.add_slide(layout)
                
                # Título
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = content.get('title', '')
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.alignment = 1  # Center
                
                # Subtítulo
                if content.get('subtitle'):
                    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
                    sub_frame = sub_box.text_frame
                    sub_para = sub_frame.paragraphs[0]
                    sub_para.text = content.get('subtitle', '')
                    sub_para.font.size = Pt(24)
                    sub_para.alignment = 1
            else:
                # Layout de contenido
                layout = prs.slide_layouts[6]  # Blank
                slide = prs.slides.add_slide(layout)
                
                # Heading
                if content.get('heading'):
                    head_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
                    head_frame = head_box.text_frame
                    head_para = head_frame.paragraphs[0]
                    head_para.text = content.get('heading', '')
                    head_para.font.size = Pt(32)
                    head_para.font.bold = True
                
                # Bullets
                bullets = content.get('bullets', [])
                if bullets:
                    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5))
                    bullet_frame = bullet_box.text_frame
                    bullet_frame.word_wrap = True
                    
                    for i, bullet in enumerate(bullets):
                        if bullet and bullet.strip():
                            if i == 0:
                                para = bullet_frame.paragraphs[0]
                            else:
                                para = bullet_frame.add_paragraph()
                            para.text = f"• {bullet}"
                            para.font.size = Pt(18)
                            para.space_after = Pt(12)
        
        # Guardar archivo
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            prs.save(tmp.name)
            tmp_path = tmp.name
        
        return FileResponse(
            tmp_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="presentacion.pptx",
            headers={"Content-Disposition": "attachment; filename=presentacion.pptx"}
        )
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al exportar PPTX: {str(e)}")

@app.post("/api/export/pdf")
async def export_pdf(request: Request):
    """
    Exporta slides a PDF usando las imágenes de preview.
    Tamaño exacto 16:9 sin márgenes.
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        from PIL import Image
        from io import BytesIO
        
        # Tamaño 16:9 en puntos (1920x1080 escalado a puntos PDF)
        # PowerPoint usa 10 pulgadas x 5.625 pulgadas para 16:9
        # En puntos: 720 x 405 (72 puntos por pulgada)
        # Usamos tamaño más grande para mejor calidad
        PAGE_WIDTH = 960  # puntos (13.33 pulgadas * 72)
        PAGE_HEIGHT = 540  # puntos (7.5 pulgadas * 72)
        
        # Leer el formulario
        content_type = request.headers.get('content-type', '')
        
        slides = []
        
        if 'multipart/form-data' in content_type:
            body = await request.body()
            
            import re
            boundary_match = re.search(r'boundary=([^;]+)', content_type)
            if boundary_match:
                boundary = boundary_match.group(1).strip('"')
                parts = body.split(f'--{boundary}'.encode())
                
                for part in parts:
                    if b'name="data"' in part:
                        content_start = part.find(b'\r\n\r\n')
                        if content_start != -1:
                            data_bytes = part[content_start + 4:]
                            if data_bytes.endswith(b'\r\n'):
                                data_bytes = data_bytes[:-2]
                            data = json.loads(data_bytes.decode('utf-8'))
                            slides = data.get('slides', [])
        else:
            body = await request.body()
            data = json.loads(body)
            slides = data.get('slides', [])
        
        logger.info(f"📄 Exportando PDF 16:9 con {len(slides)} slides")
        
        # Crear PDF con tamaño exacto 16:9
        pdf_path = tempfile.mktemp(suffix='.pdf')
        c = canvas.Canvas(pdf_path, pagesize=(PAGE_WIDTH, PAGE_HEIGHT))
        
        for idx, slide in enumerate(slides):
            preview = slide.get('preview')
            
            if preview and preview.startswith('data:image'):
                try:
                    # Decodificar imagen base64
                    header, base64_data = preview.split(',', 1)
                    image_data = base64.b64decode(base64_data)
                    
                    # Abrir con Pillow
                    img = Image.open(BytesIO(image_data))
                    
                    # Convertir a RGB si es necesario
                    if img.mode in ('RGBA', 'P'):
                        img = img.convert('RGB')
                    
                    # Guardar imagen temporal
                    img_buffer = BytesIO()
                    img.save(img_buffer, format='PNG')
                    img_buffer.seek(0)
                    
                    # Dibujar imagen ocupando TODA la página (sin márgenes)
                    c.drawImage(ImageReader(img_buffer), 0, 0, PAGE_WIDTH, PAGE_HEIGHT)
                    
                    logger.info(f"   ✅ Slide {idx + 1} agregado al PDF")
                    
                except Exception as e:
                    logger.warning(f"   ⚠️ Error procesando imagen slide {idx + 1}: {e}")
                    # Crear página con texto de fallback
                    c.setFont("Helvetica-Bold", 24)
                    c.drawCentredString(PAGE_WIDTH / 2, PAGE_HEIGHT / 2, f"Slide {idx + 1}")
            else:
                # Sin preview - crear página con contenido de texto
                c.setFont("Helvetica-Bold", 32)
                content = slide.get('content', {})
                title = content.get('title') or content.get('heading', f'Slide {idx + 1}')
                c.drawCentredString(PAGE_WIDTH / 2, PAGE_HEIGHT - 80, title)
                
                if content.get('subtitle'):
                    c.setFont("Helvetica", 20)
                    c.drawCentredString(PAGE_WIDTH / 2, PAGE_HEIGHT - 120, content['subtitle'])
                
                bullets = content.get('bullets', [])
                if bullets:
                    c.setFont("Helvetica", 14)
                    y_pos = PAGE_HEIGHT - 180
                    for bullet in bullets:
                        if bullet and bullet.strip():
                            c.drawString(80, y_pos, f"• {bullet}")
                            y_pos -= 22
                
                logger.info(f"   📝 Slide {idx + 1} (texto) agregado al PDF")
            
            # Nueva página si no es el último slide
            if idx < len(slides) - 1:
                c.showPage()
        
        c.save()
        
        logger.info(f"✅ PDF 16:9 generado: {pdf_path}")
        
        return FileResponse(
            pdf_path,
            media_type="application/pdf",
            filename="presentacion.pdf",
            headers={"Content-Disposition": "attachment; filename=presentacion.pdf"}
        )
    
    except Exception as e:
        logger.error(f"❌ Error exportando PDF: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al exportar PDF: {str(e)}")

@app.post("/api/extract-content")
async def extract_content(file: UploadFile = File(...)):
    """
    Extrae solo el contenido de texto de un PPTX (sin diseño)
    Útil para importar contenido generado por IA
    """
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    try:
        # Guardar archivo temporalmente
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name
        
        print(f"📄 Extrayendo contenido de: {tmp_path}")
        
        # Abrir presentación
        prs = Presentation(tmp_path)
        
        extracted_slides = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = {
                "slideNumber": slide_idx + 1,
                "type": "content",
                "texts": []
            }
            
            # Extraer todo el texto del slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # Detectar tipo de texto
                        text_type = "body"
                        if shape.is_placeholder:
                            placeholder_type = shape.placeholder_format.type
                            if placeholder_type == 1:
                                text_type = "title"
                            elif placeholder_type == 2:
                                text_type = "subtitle"
                        
                        # Si el texto tiene bullets, separar
                        if '\n' in text and text_type == "body":
                            bullets = [line.strip() for line in text.split('\n') if line.strip()]
                            slide_content["texts"].append({
                                "type": "bullets",
                                "content": bullets
                            })
                        else:
                            slide_content["texts"].append({
                                "type": text_type,
                                "content": text
                            })
            
            # Detectar tipo de slide
            if slide_idx == 0 or any(t["type"] == "title" for t in slide_content["texts"]):
                slide_content["type"] = "title"
            
            extracted_slides.append(slide_content)
        
        # Limpiar archivo temporal
        os.unlink(tmp_path)
        
        print(f"✅ Contenido extraído: {len(extracted_slides)} slides")
        
        return {
            "success": True,
            "fileName": file.filename,
            "slideCount": len(extracted_slides),
            "slides": extracted_slides
        }
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al extraer contenido: {str(e)}")

# ============================================
# ENDPOINT DE ANÁLISIS DE FUENTES
# ============================================

@app.post("/api/analyze-fonts")
async def analyze_fonts_endpoint(file: UploadFile = File(...)):
    """
    Analiza las fuentes usadas en un PPTX y detecta cuáles faltan.
    
    Returns:
        - fonts: Lista de todas las fuentes detectadas
        - missingFonts: Fuentes que no están en el sistema
        - availableOnline: Fuentes disponibles en Google Fonts
        - googleFontsLink: URL para cargar las fuentes de Google
        - warnings: Advertencias sobre fuentes faltantes
    """
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    try:
        # Guardar archivo temporalmente
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name
        
        logger.info(f"🔤 Analizando fuentes de: {file.filename}")
        
        # Analizar fuentes
        font_analysis = analyze_pptx_fonts(tmp_path)
        
        # Limpiar archivo temporal
        os.unlink(tmp_path)
        
        logger.info(f"✅ Análisis de fuentes completado: {font_analysis['summary']}")
        
        return {
            "success": True,
            "fileName": file.filename,
            **font_analysis
        }
    
    except Exception as e:
        logger.error(f"❌ Error analizando fuentes: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al analizar fuentes: {str(e)}")


# ============================================
# ENDPOINTS DE TEMPLATE AUTO-MAPPING
# ============================================

# Gemini Vision API configuration
GEMINI_API_KEY = os.environ.get('VITE_GEMINI_API_KEY', os.environ.get('GEMINI_API_KEY', ''))
GEMINI_MODEL = os.environ.get('VITE_GEMINI_MODEL', 'gemini-1.5-flash')
GEMINI_API_URL = 'https://generativelanguage.googleapis.com/v1beta/models'

# Technical prompt for template layout analysis
TEMPLATE_ANALYSIS_PROMPT = """Actúa como un analizador experto de XML de PowerPoint.

Tarea: Analiza la imagen adjunta e identifica todos los contenedores de contenido (placeholders).

Instrucciones de Extracción:
1. Identifica el propósito de cada área: TITLE, SUBTITLE, BODY, FOOTER, IMAGE_HOLDER, o CHART_AREA
2. Devuelve las coordenadas en formato normalizado (0-1000) para top, left, width, height
3. Identifica atributos visuales: font_color (hex), text_alignment (left, center, right), y background_color
4. Estima el "Z-index" (orden de capas) si hay elementos superpuestos

Formato de salida: JSON puro. No agregues explicaciones.

{
  "slide_metadata": { "aspect_ratio": "16:9" },
  "elements": [
    {
      "id": "element_1",
      "type": "TITLE",
      "coordinates": {"top": 50, "left": 100, "width": 800, "height": 100},
      "style": {"color": "#2C3E50", "align": "center"},
      "confidence": 0.95
    }
  ]
}"""

# Valid element types
VALID_ELEMENT_TYPES = ['TITLE', 'SUBTITLE', 'BODY', 'FOOTER', 'IMAGE_HOLDER', 'CHART_AREA', 'UNKNOWN']


async def call_gemini_vision_api(base64_data: str) -> dict:
    """
    Make a single Gemini Vision API call.
    
    Args:
        base64_data: Base64 encoded image data (without prefix)
        
    Returns:
        Raw JSON response from Gemini
        
    Raises:
        Exception if API call fails
    """
    if not GEMINI_API_KEY:
        raise HTTPException(status_code=500, detail="GEMINI_API_KEY not configured")
    
    url = f"{GEMINI_API_URL}/{GEMINI_MODEL}:generateContent?key={GEMINI_API_KEY}"
    
    payload = {
        "contents": [
            {
                "parts": [
                    {"text": TEMPLATE_ANALYSIS_PROMPT},
                    {
                        "inline_data": {
                            "mime_type": "image/png",
                            "data": base64_data
                        }
                    }
                ]
            }
        ],
        "generationConfig": {
            "temperature": 0.2,
            "maxOutputTokens": 4096,
            "responseMimeType": "application/json"
        }
    }
    
    async with httpx.AsyncClient(timeout=60.0) as client:
        response = await client.post(url, json=payload)
        
        if response.status_code != 200:
            error_data = response.json() if response.content else {}
            error_msg = error_data.get('error', {}).get('message', 'Unknown error')
            raise Exception(f"Gemini API Error: {response.status_code} - {error_msg}")
        
        data = response.json()
        text_content = data.get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', '{}')
        
        # Extract JSON from response
        import re
        json_match = re.search(r'\{[\s\S]*\}', text_content)
        if not json_match:
            raise Exception('No valid JSON found in Gemini response')
        
        return json.loads(json_match.group(0))


def normalize_coordinate(value: Any) -> int:
    """Normalize a coordinate value to the range [0, 1000]"""
    if not isinstance(value, (int, float)):
        return 0
    return max(0, min(1000, round(value)))


def validate_element_type(element_type: str) -> str:
    """Validate and normalize element type"""
    if not isinstance(element_type, str):
        return 'UNKNOWN'
    upper_type = element_type.upper().strip()
    return upper_type if upper_type in VALID_ELEMENT_TYPES else 'UNKNOWN'


def parse_vision_response(raw_response: dict) -> dict:
    """
    Parse and validate Gemini Vision response for template analysis.
    Normalizes coordinates to range [0, 1000] and validates element types.
    """
    if not raw_response or not isinstance(raw_response, dict):
        return {
            'slide_metadata': {'aspect_ratio': '16:9'},
            'elements': []
        }
    
    slide_metadata = {
        'aspect_ratio': raw_response.get('slide_metadata', {}).get('aspect_ratio', '16:9')
    }
    
    raw_elements = raw_response.get('elements', [])
    if not isinstance(raw_elements, list):
        raw_elements = []
    
    elements = []
    for idx, element in enumerate(raw_elements):
        if not isinstance(element, dict):
            continue
            
        coords = element.get('coordinates', {})
        normalized_coords = {
            'top': normalize_coordinate(coords.get('top')),
            'left': normalize_coordinate(coords.get('left')),
            'width': normalize_coordinate(coords.get('width')),
            'height': normalize_coordinate(coords.get('height'))
        }
        
        validated_type = validate_element_type(element.get('type', 'UNKNOWN'))
        
        style = element.get('style', {})
        normalized_style = {
            'color': style.get('color', '#000000') if isinstance(style.get('color'), str) else '#000000',
            'align': style.get('align', 'left') if style.get('align') in ['left', 'center', 'right'] else 'left',
            'backgroundColor': style.get('backgroundColor') if isinstance(style.get('backgroundColor'), str) else None
        }
        
        confidence = element.get('confidence', 0.5)
        if not isinstance(confidence, (int, float)):
            confidence = 0.5
        confidence = max(0, min(1, confidence))
        
        elements.append({
            'id': element.get('id', f'element_{idx + 1}'),
            'type': validated_type,
            'coordinates': normalized_coords,
            'style': normalized_style,
            'confidence': confidence,
            'shapeId': None  # Will be populated after shape matching
        })
    
    return {
        'slide_metadata': slide_metadata,
        'elements': elements
    }


async def analyze_with_retry(image_base64: str, max_retries: int = 2) -> dict:
    """
    Analyze template layout with retry logic and exponential backoff.
    Retries up to 2 times with backoff: 1s, 2s
    
    Args:
        image_base64: Base64 encoded image of the slide
        max_retries: Maximum number of retries (default: 2)
        
    Returns:
        Parsed vision response with elements
        
    Raises:
        HTTPException if all retries fail
    """
    # Remove data URL prefix if present
    base64_data = image_base64
    if image_base64.startswith('data:'):
        base64_data = image_base64.split(',', 1)[1] if ',' in image_base64 else image_base64
    
    last_error = None
    
    for attempt in range(max_retries + 1):
        try:
            if attempt > 0:
                backoff_ms = 1000 * attempt  # 1s, 2s
                logger.info(f"⏳ Retry attempt {attempt}/{max_retries} after {backoff_ms}ms backoff...")
                await asyncio.sleep(backoff_ms / 1000)
            
            raw_result = await call_gemini_vision_api(base64_data)
            logger.info('📄 Raw Gemini response received')
            
            parsed = parse_vision_response(raw_result)
            
            # Validate that we got at least some elements
            if len(parsed['elements']) == 0 and attempt < max_retries:
                logger.warning('⚠️ Empty elements array, retrying...')
                last_error = Exception('Gemini returned empty elements array')
                continue
            
            logger.info(f"✅ Template analysis completed with {len(parsed['elements'])} elements")
            return parsed
            
        except Exception as e:
            last_error = e
            logger.error(f"❌ Attempt {attempt + 1}/{max_retries + 1} failed: {str(e)}")
            
            if attempt == max_retries:
                break
    
    raise HTTPException(
        status_code=500, 
        detail=f"Gemini Vision analysis failed after {max_retries + 1} attempts: {str(last_error)}"
    )


def convert_slide_to_image(pptx_path: str, slide_index: int = 0) -> str:
    """
    Convert a specific slide to a base64 PNG image.
    Uses LibreOffice or custom renderer.
    
    Args:
        pptx_path: Path to the PPTX file
        slide_index: Index of the slide to convert (default: 0 for first slide)
        
    Returns:
        Base64 encoded PNG image with data URL prefix
    """
    try:
        # Try to use the existing conversion functions
        from pptx_to_images import convert_pptx_to_images
        
        images = convert_pptx_to_images(pptx_path)
        if images and len(images) > slide_index:
            return images[slide_index]
    except Exception as e:
        logger.warning(f"LibreOffice conversion failed: {e}")
    
    # Fallback to custom renderer
    try:
        from pptx_to_images_custom import convert_pptx_to_images_custom
        
        images = convert_pptx_to_images_custom(pptx_path)
        if images and len(images) > slide_index:
            return images[slide_index]
    except Exception as e:
        logger.warning(f"Custom renderer failed: {e}")
    
    # Final fallback: generate placeholder
    from pptx_to_images import generate_placeholder_image
    return generate_placeholder_image(slide_index + 1)


@app.post("/api/analyze-template")
async def analyze_template(file: UploadFile = File(...)):
    """
    Analyze a PPTX template and return mapping of elements.
    
    1. Verifies cache by hash
    2. If not cached, converts to image and calls Gemini Vision
    3. Links detections with real shapes using ShapeMatcher
    4. Saves to cache and returns
    
    Requirements: 1.1, 1.2, 4.3, 6.5
    """
    operation_context = {
        "filename": file.filename,
        "content_type": file.content_type
    }
    log_operation_start("analyze_template", operation_context)
    
    if not file.filename.endswith(('.pptx', '.ppt')):
        log_error_with_context(
            ValueError("Invalid file type"),
            ErrorCategory.FILE_VALIDATION,
            "validate_file_type",
            {"filename": file.filename, "expected": [".pptx", ".ppt"]},
            include_traceback=False
        )
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    tmp_path = None
    
    try:
        # Read file content
        file_content = await file.read()
        operation_context["file_size"] = len(file_content)
        
        # Generate hash for cache lookup
        template_hash = mapping_cache.generate_template_hash(file_content)
        operation_context["template_hash"] = template_hash[:16] + "..."
        logger.info(f"📄 Template hash: {template_hash[:16]}...")
        
        # Check cache first (Requirement 4.3)
        cached_mapping = mapping_cache.get_cached_mapping(template_hash)
        if cached_mapping:
            log_operation_success("analyze_template", f"Loaded from cache - {len(cached_mapping['elements'])} elements")
            return {
                "success": True,
                "templateHash": template_hash,
                "elements": cached_mapping['elements'],
                "shapeMapping": cached_mapping['shapeMapping'],
                "source": "cache",
                "message": f"Template reconocido - {len(cached_mapping['elements'])} elementos cargados desde caché"
            }
        
        # Save file temporarily for analysis
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name
        
        logger.info(f"📄 Archivo guardado temporalmente: {tmp_path}")
        
        # Open presentation to get slide dimensions and shapes
        try:
            prs = Presentation(tmp_path)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            operation_context["slide_dimensions"] = f"{slide_width}x{slide_height}"
        except Exception as e:
            log_error_with_context(
                e,
                ErrorCategory.FILE_CONVERSION,
                "open_presentation",
                operation_context
            )
            raise HTTPException(
                status_code=400, 
                detail="No se pudo abrir el archivo PPTX. Verifica que el archivo no esté corrupto."
            )
        
        # Get shapes from first slide (or slide master)
        shapes = []
        if prs.slides:
            shapes = list(prs.slides[0].shapes)
        operation_context["shape_count"] = len(shapes)
        
        logger.info(f"📐 Slide dimensions: {slide_width} x {slide_height} EMUs")
        logger.info(f"🔷 Found {len(shapes)} shapes in first slide")
        
        # Convert first slide to image (Requirement 1.1)
        logger.info("🎨 Converting slide to image...")
        try:
            slide_image = convert_slide_to_image(tmp_path, 0)
            if not slide_image:
                raise ValueError("Image conversion returned empty result")
        except Exception as e:
            log_error_with_context(
                e,
                ErrorCategory.FILE_CONVERSION,
                "convert_slide_to_image",
                operation_context
            )
            raise HTTPException(
                status_code=500, 
                detail="Error de conversión: No se pudo convertir el slide a imagen. Verifica que el archivo sea un PPTX válido."
            )
        
        # Call Gemini Vision with retry (Requirement 1.2)
        logger.info("🔍 Calling Gemini Vision API...")
        try:
            vision_result = await analyze_with_retry(slide_image)
            operation_context["detected_elements"] = len(vision_result.get('elements', []))
        except HTTPException as e:
            log_error_with_context(
                Exception(e.detail),
                ErrorCategory.GEMINI_API,
                "analyze_with_retry",
                operation_context
            )
            raise
        except Exception as e:
            log_error_with_context(
                e,
                ErrorCategory.GEMINI_API,
                "analyze_with_retry",
                operation_context
            )
            raise HTTPException(
                status_code=500, 
                detail=f"Error de análisis visual: {str(e)}"
            )
        
        # Match detections to shapes using ShapeMatcher
        logger.info("🔗 Matching detections to shapes...")
        try:
            matcher = ShapeMatcher(slide_width, slide_height)
            
            # Prepare detections for matching
            detections = vision_result['elements']
            
            # Match detections to shapes
            shape_mapping = matcher.match_detections_to_shapes(shapes, detections)
            
            # Update elements with shape IDs
            for element in detections:
                element_type = element['type']
                if element_type in shape_mapping:
                    element['shapeId'] = shape_mapping[element_type]
            
            # Also classify shapes that weren't matched by vision
            for shape in shapes:
                shape_id = getattr(shape, 'shape_id', None)
                if shape_id and shape_id not in shape_mapping.values():
                    # Classify by geometry or metadata
                    classified_type = matcher.classify_shape(shape)
                    if classified_type != 'UNKNOWN' and classified_type not in shape_mapping:
                        shape_mapping[classified_type] = shape_id
            
            operation_context["shape_mapping"] = shape_mapping
            logger.info(f"✅ Shape mapping complete: {shape_mapping}")
            
        except Exception as e:
            log_error_with_context(
                e,
                ErrorCategory.SHAPE_MATCHING,
                "match_detections_to_shapes",
                operation_context
            )
            # Continue with empty mapping - don't fail the whole operation
            logger.warning("⚠️ Shape matching failed, continuing with empty mapping")
            shape_mapping = {}
        
        # Save to cache (Requirement 4.1)
        try:
            mapping_to_save = {
                'elements': detections,
                'shapeMapping': shape_mapping
            }
            mapping_cache.save_mapping(
                template_hash=template_hash,
                mapping=mapping_to_save,
                template_name=file.filename,
                file_path=None,  # We don't persist the file
                thumbnail_url=slide_image[:100] + '...' if slide_image else None  # Just store a reference
            )
            logger.info(f"💾 Mapping saved to cache")
        except Exception as e:
            log_error_with_context(
                e,
                ErrorCategory.CACHE_OPERATION,
                "save_mapping",
                operation_context
            )
            # Don't fail the operation if cache save fails
            logger.warning("⚠️ Failed to save mapping to cache, continuing...")
        
        log_operation_success("analyze_template", f"{len(detections)} elements detected")
        
        return {
            "success": True,
            "templateHash": template_hash,
            "elements": detections,
            "shapeMapping": shape_mapping,
            "source": "vision",
            "message": f"Análisis completado - {len(detections)} elementos detectados"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        log_error_with_context(
            e,
            ErrorCategory.UNKNOWN,
            "analyze_template",
            operation_context
        )
        raise HTTPException(status_code=500, detail=f"Error al analizar template: {str(e)}")
    
    finally:
        # Clean up temporary file (Requirement 6.4 - preserve original)
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
                logger.info(f"🗑️ Cleaned up temp file: {tmp_path}")
            except Exception as e:
                logger.warning(f"Could not delete temp file: {e}")


@app.post("/api/update-mapping")
async def update_mapping(request: UpdateMappingRequest):
    """
    Update element type in cache (manual user correction).
    
    When a user corrects an incorrect classification,
    this correction is persisted for future uses of the same template.
    
    Requirements: 2.4, 2.5, 6.5
    """
    operation_context = {
        "template_hash": request.template_hash[:16] + "...",
        "element_id": request.element_id,
        "new_type": request.new_type
    }
    log_operation_start("update_mapping", operation_context)
    
    try:
        # Validate new_type
        validated_type = validate_element_type(request.new_type)
        if validated_type == 'UNKNOWN' and request.new_type.upper() != 'UNKNOWN':
            log_error_with_context(
                ValueError(f"Invalid element type: {request.new_type}"),
                ErrorCategory.FILE_VALIDATION,
                "validate_element_type",
                operation_context,
                include_traceback=False
            )
            raise HTTPException(
                status_code=400, 
                detail=f"Invalid element type: {request.new_type}. Valid types: {', '.join(VALID_ELEMENT_TYPES)}"
            )
        
        # Update in cache
        success = mapping_cache.update_element_type(
            template_hash=request.template_hash,
            element_id=request.element_id,
            new_type=validated_type
        )
        
        if not success:
            log_error_with_context(
                ValueError("Element not found"),
                ErrorCategory.CACHE_OPERATION,
                "update_element_type",
                operation_context,
                include_traceback=False
            )
            raise HTTPException(
                status_code=404, 
                detail=f"Element {request.element_id} not found for template {request.template_hash[:16]}..."
            )
        
        log_operation_success("update_mapping", f"Element {request.element_id} updated to {validated_type}")
        
        return {
            "success": True,
            "message": f"Elemento actualizado a tipo {validated_type}"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        log_error_with_context(
            e,
            ErrorCategory.UNKNOWN,
            "update_mapping",
            operation_context
        )
        raise HTTPException(status_code=500, detail=f"Error al actualizar mapping: {str(e)}")


@app.get("/api/templates")
async def list_templates():
    """
    List all cached templates.
    """
    try:
        templates = mapping_cache.get_all_templates()
        return {
            "success": True,
            "templates": templates,
            "count": len(templates)
        }
    except Exception as e:
        logger.error(f"❌ Error listing templates: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al listar templates: {str(e)}")


@app.get("/api/template/{template_hash}")
async def get_template_mapping(template_hash: str):
    """
    Get cached mapping for a specific template by hash.
    """
    try:
        cached_mapping = mapping_cache.get_cached_mapping(template_hash)
        
        if not cached_mapping:
            raise HTTPException(status_code=404, detail="Template no encontrado en caché")
        
        return {
            "success": True,
            "mapping": cached_mapping
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error getting template: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al obtener template: {str(e)}")


@app.delete("/api/template/{template_hash}")
async def delete_template_mapping(template_hash: str):
    """
    Delete a cached template mapping.
    """
    try:
        success = mapping_cache.delete_mapping(template_hash)
        
        if not success:
            raise HTTPException(status_code=404, detail="Template no encontrado en caché")
        
        return {
            "success": True,
            "message": "Template eliminado de caché"
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error deleting template: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al eliminar template: {str(e)}")


# ============================================
# ENDPOINTS DE COLABORACIÓN
# ============================================

@app.post("/api/presentations/create")
async def create_shared_presentation(data: Dict[str, Any] = Body(...)):
    """
    Crear presentación compartida
    """
    try:
        owner = data.get('owner', 'anonymous')
        title = data.get('title', 'Presentación sin título')
        template_data = data.get('templateData', {})
        slides_data = data.get('slidesData', [])
        extracted_assets = data.get('extractedAssets')
        permissions = data.get('permissions')
        
        presentation_id = db.create_presentation(
            owner=owner,
            title=title,
            template_data=template_data,
            slides_data=slides_data,
            extracted_assets=extracted_assets,
            permissions=permissions
        )
        
        return {
            "success": True,
            "presentationId": presentation_id,
            "shareUrl": f"/editor/{presentation_id}"
        }
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al crear presentación: {str(e)}")

@app.get("/api/presentations/{presentation_id}")
async def get_shared_presentation(presentation_id: str):
    """
    Obtener presentación compartida
    """
    try:
        presentation = db.get_presentation(presentation_id)
        
        if not presentation:
            raise HTTPException(status_code=404, detail="Presentación no encontrada")
        
        return {
            "success": True,
            "presentation": presentation
        }
    
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al obtener presentación: {str(e)}")

@app.put("/api/presentations/{presentation_id}")
async def update_shared_presentation(
    presentation_id: str,
    data: Dict[str, Any] = Body(...)
):
    """
    Actualizar presentación compartida
    """
    try:
        user = data.get('user', 'anonymous')
        slides_data = data.get('slidesData', [])
        
        # Verificar permisos
        if not db.check_permission(presentation_id, user, 'edit'):
            raise HTTPException(status_code=403, detail="No tienes permiso para editar")
        
        success = db.update_presentation(presentation_id, slides_data, user)
        
        if not success:
            raise HTTPException(status_code=404, detail="Presentación no encontrada")
        
        # Notificar a otros usuarios conectados
        await manager.broadcast(presentation_id, {
            "type": "presentation_updated",
            "user": user,
            "slidesData": slides_data
        })
        
        return {
            "success": True,
            "message": "Presentación actualizada"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al actualizar presentación: {str(e)}")

@app.put("/api/presentations/{presentation_id}/permissions")
async def update_presentation_permissions(
    presentation_id: str,
    data: Dict[str, Any] = Body(...)
):
    """
    Actualizar permisos de presentación
    """
    try:
        user = data.get('user', 'anonymous')
        permissions = data.get('permissions', {})
        
        success = db.update_permissions(presentation_id, permissions, user)
        
        if not success:
            raise HTTPException(status_code=403, detail="Solo el propietario puede cambiar permisos")
        
        return {
            "success": True,
            "message": "Permisos actualizados"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al actualizar permisos: {str(e)}")

@app.websocket("/ws/{presentation_id}")
async def websocket_endpoint(websocket: WebSocket, presentation_id: str):
    """
    WebSocket para colaboración en tiempo real
    """
    await manager.connect(websocket, presentation_id)
    
    try:
        while True:
            data = await websocket.receive_json()
            
            message_type = data.get('type')
            user = data.get('user', 'anonymous')
            
            if message_type == 'slide_update':
                # Actualizar slide en DB
                slide_index = data.get('slideIndex')
                slide_data = data.get('slideData')
                
                db.update_slide(presentation_id, slide_index, slide_data, user)
                
                # Broadcast a otros usuarios
                await manager.broadcast(presentation_id, {
                    "type": "slide_updated",
                    "user": user,
                    "slideIndex": slide_index,
                    "slideData": slide_data
                }, exclude=websocket)
            
            elif message_type == 'cursor_move':
                # Broadcast posición del cursor
                await manager.broadcast(presentation_id, {
                    "type": "cursor_moved",
                    "user": user,
                    "position": data.get('position')
                }, exclude=websocket)
            
            elif message_type == 'user_joined':
                # Notificar que usuario se unió
                await manager.broadcast(presentation_id, {
                    "type": "user_joined",
                    "user": user
                }, exclude=websocket)
    
    except WebSocketDisconnect:
        manager.disconnect(websocket, presentation_id)
        await manager.broadcast(presentation_id, {
            "type": "user_left",
            "user": "unknown"
        })
    except Exception as e:
        print(f"❌ Error en WebSocket: {e}")
        manager.disconnect(websocket, presentation_id)

if __name__ == "__main__":
    # Configurar límite de tamaño de archivo grande (50MB)
    uvicorn.run(
        app, 
        host="0.0.0.0", 
        port=8000,
        # Aumentar límites para archivos grandes
        limit_concurrency=100,
        timeout_keep_alive=30
    )

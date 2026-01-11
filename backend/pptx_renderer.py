"""
Renderizador personalizado de PPTX que integra correctamente logos con el fondo
"""
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
import base64
from typing import List

def render_slide_to_image(slide, prs, slide_number: int) -> str:
    """
    Renderiza un slide completo a imagen, integrando correctamente todos los elementos
    """
    # Dimensiones del slide en EMUs
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    
    # Convertir a pixels (96 DPI)
    # 1 EMU = 1/914400 pulgadas
    # 1 pulgada = 96 pixels (para pantalla)
    dpi = 96
    width_px = int(slide_width_emu / 914400 * dpi * 3)  # 3x para mejor calidad
    height_px = int(slide_height_emu / 914400 * dpi * 3)
    
    # Crear imagen base
    img = Image.new('RGB', (width_px, height_px), color='white')
    draw = ImageDraw.Draw(img)
    
    # 1. Renderizar fondo
    try:
        bg_color = extract_background_color(slide)
        if bg_color and bg_color != '#FFFFFF':
            img = Image.new('RGB', (width_px, height_px), color=bg_color)
            draw = ImageDraw.Draw(img)
            print(f"   🎨 Fondo aplicado: {bg_color}")
    except Exception as e:
        print(f"   ⚠️ Error aplicando fondo: {e}")
    
    # 2. Renderizar shapes (imágenes, formas, etc.)
    for shape in slide.shapes:
        try:
            # Calcular posición en pixels
            left_px = int(shape.left / 914400 * dpi * 3)
            top_px = int(shape.top / 914400 * dpi * 3)
            width_shape_px = int(shape.width / 914400 * dpi * 3)
            height_shape_px = int(shape.height / 914400 * dpi * 3)
            
            # Renderizar imágenes
            if shape.shape_type == 13:  # PICTURE
                try:
                    image_bytes = shape.image.blob
                    shape_img = Image.open(BytesIO(image_bytes))
                    
                    # Redimensionar a la posición del shape
                    shape_img = shape_img.resize((width_shape_px, height_shape_px), Image.LANCZOS)
                    
                    # Si la imagen tiene transparencia, usarla
                    if shape_img.mode in ('RGBA', 'LA'):
                        img.paste(shape_img, (left_px, top_px), shape_img)
                    else:
                        # Si no tiene transparencia pero tiene fondo blanco, intentar removerlo
                        if has_white_background(shape_img):
                            shape_img = remove_white_bg(shape_img, bg_color if bg_color else '#FFFFFF')
                        img.paste(shape_img, (left_px, top_px))
                    
                    print(f"   ✅ Imagen renderizada en ({left_px}, {top_px})")
                except Exception as e:
                    print(f"   ⚠️ Error renderizando imagen: {e}")
            
            # Renderizar texto (simplificado)
            elif shape.has_text_frame:
                try:
                    text = shape.text_frame.text
                    if text.strip():
                        # Usar fuente por defecto
                        try:
                            font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 40)
                        except:
                            font = ImageFont.load_default()
                        
                        draw.text((left_px + 10, top_px + 10), text, fill='black', font=font)
                except Exception as e:
                    print(f"   ⚠️ Error renderizando texto: {e}")
        
        except Exception as e:
            print(f"   ⚠️ Error procesando shape: {e}")
    
    # Convertir a base64
    buffer = BytesIO()
    img.save(buffer, format='PNG')
    img_str = base64.b64encode(buffer.getvalue()).decode()
    
    return f"data:image/png;base64,{img_str}"


def extract_background_color(slide) -> str:
    """Extrae el color de fondo del slide"""
    try:
        from lxml import etree
        
        slide_part = slide.part
        slide_xml = slide_part.blob
        root = etree.fromstring(slide_xml)
        
        namespaces = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        
        # Buscar fondo
        bg_element = root.find('.//p:cSld/p:bg', namespaces)
        if bg_element is not None:
            solid_fill = bg_element.find('.//a:solidFill/a:srgbClr', namespaces)
            if solid_fill is not None:
                color_val = solid_fill.get('val')
                if color_val:
                    return f"#{color_val}"
        
        # Buscar en layout
        layout_part = slide.slide_layout.part
        layout_xml = layout_part.blob
        layout_root = etree.fromstring(layout_xml)
        
        bg_element = layout_root.find('.//p:cSld/p:bg', namespaces)
        if bg_element is not None:
            solid_fill = bg_element.find('.//a:solidFill/a:srgbClr', namespaces)
            if solid_fill is not None:
                color_val = solid_fill.get('val')
                if color_val:
                    return f"#{color_val}"
    
    except Exception as e:
        print(f"   ⚠️ Error extrayendo color de fondo: {e}")
    
    return '#FFFFFF'


def has_white_background(img: Image.Image) -> bool:
    """Detecta si una imagen tiene fondo blanco"""
    # Muestrear esquinas
    corners = [
        img.getpixel((0, 0)),
        img.getpixel((img.width-1, 0)),
        img.getpixel((0, img.height-1)),
        img.getpixel((img.width-1, img.height-1))
    ]
    
    white_count = 0
    for pixel in corners:
        if img.mode == 'RGB':
            r, g, b = pixel
            if r > 240 and g > 240 and b > 240:
                white_count += 1
        elif img.mode == 'RGBA':
            r, g, b, a = pixel
            if r > 240 and g > 240 and b > 240 and a > 200:
                white_count += 1
    
    return white_count >= 3  # Al menos 3 esquinas blancas


def remove_white_bg(img: Image.Image, bg_color: str) -> Image.Image:
    """Remueve fondo blanco y lo reemplaza con el color del slide"""
    img = img.convert('RGBA')
    data = img.getdata()
    
    # Convertir color hex a RGB
    bg_color = bg_color.lstrip('#')
    bg_r = int(bg_color[0:2], 16)
    bg_g = int(bg_color[2:4], 16)
    bg_b = int(bg_color[4:6], 16)
    
    new_data = []
    for item in data:
        r, g, b, a = item
        # Si es blanco (o casi blanco)
        if r > 240 and g > 240 and b > 240:
            new_data.append((bg_r, bg_g, bg_b, a))
        else:
            new_data.append(item)
    
    img.putdata(new_data)
    return img


def render_pptx_to_images(pptx_path: str) -> List[str]:
    """
    Renderiza todas las slides de un PPTX a imágenes
    """
    print(f"\n🎨 Renderizando PPTX con renderizador personalizado...")
    
    prs = Presentation(pptx_path)
    images = []
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n📄 Renderizando slide {i}...")
        img_base64 = render_slide_to_image(slide, prs, i)
        images.append(img_base64)
        print(f"✅ Slide {i} renderizado")
    
    print(f"\n✅ {len(images)} slides renderizados\n")
    return images

"""
Analizar la imagen del logo extraída del PPTX para ver si tiene fondo blanco o transparente
"""
from pptx import Presentation
from PIL import Image
from io import BytesIO
import sys

def analyze_logo(pptx_path):
    """Analiza las imágenes del PPTX"""
    prs = Presentation(pptx_path)
    
    print(f"\n{'='*80}")
    print(f"ANÁLISIS DE IMÁGENES EN EL PPTX")
    print(f"{'='*80}\n")
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"\n--- SLIDE {slide_idx} ---")
        
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 6:  # PICTURE
                image_count += 1
                try:
                    image = shape.image
                    image_bytes = image.blob
                    content_type = image.content_type
                    
                    print(f"\n📷 Imagen {image_count}:")
                    print(f"   Shape ID: {shape.shape_id}")
                    print(f"   Formato: {content_type}")
                    print(f"   Tamaño: {len(image_bytes)} bytes")
                    
                    # Analizar la imagen con Pillow
                    img = Image.open(BytesIO(image_bytes))
                    print(f"   Modo: {img.mode}")
                    print(f"   Dimensiones: {img.size}")
                    
                    if img.mode == 'RGBA':
                        print(f"   ✅ Tiene canal alpha (transparencia)")
                        
                        # Analizar si realmente usa transparencia
                        alpha_channel = img.split()[3]
                        alpha_values = list(alpha_channel.getdata())
                        
                        min_alpha = min(alpha_values)
                        max_alpha = max(alpha_values)
                        avg_alpha = sum(alpha_values) / len(alpha_values)
                        
                        print(f"   Alpha min: {min_alpha}, max: {max_alpha}, promedio: {avg_alpha:.1f}")
                        
                        if min_alpha == 255 and max_alpha == 255:
                            print(f"   ⚠️ Todos los pixels son opacos (alpha=255) - NO usa transparencia")
                        else:
                            print(f"   ✅ Usa transparencia real")
                    
                    # Analizar colores de fondo (esquinas)
                    print(f"\n   🎨 Análisis de colores en esquinas:")
                    corners = [
                        (0, 0, "Superior izquierda"),
                        (img.width-1, 0, "Superior derecha"),
                        (0, img.height-1, "Inferior izquierda"),
                        (img.width-1, img.height-1, "Inferior derecha")
                    ]
                    
                    for x, y, name in corners:
                        pixel = img.getpixel((x, y))
                        if img.mode == 'RGBA':
                            r, g, b, a = pixel
                            hex_color = f"#{r:02x}{g:02x}{b:02x}"
                            print(f"      {name}: {hex_color} (alpha={a})")
                            
                            # Detectar si es blanco
                            if r > 240 and g > 240 and b > 240:
                                print(f"         ⚠️ BLANCO detectado")
                        else:
                            r, g, b = pixel
                            hex_color = f"#{r:02x}{g:02x}{b:02x}"
                            print(f"      {name}: {hex_color}")
                    
                    # Analizar centro de la imagen
                    center_x, center_y = img.width // 2, img.height // 2
                    center_pixel = img.getpixel((center_x, center_y))
                    if img.mode == 'RGBA':
                        r, g, b, a = center_pixel
                        print(f"\n   🎯 Centro: #{r:02x}{g:02x}{b:02x} (alpha={a})")
                    else:
                        r, g, b = center_pixel
                        print(f"\n   🎯 Centro: #{r:02x}{g:02x}{b:02x}")
                    
                except Exception as e:
                    print(f"   ❌ Error: {e}")
                    import traceback
                    traceback.print_exc()
        
        if image_count == 0:
            print("   (No hay imágenes en este slide)")
    
    print(f"\n{'='*80}\n")

if __name__ == "__main__":
    import glob
    import os
    
    # Buscar archivo más reciente en temp
    temp_files = glob.glob(r"C:\Users\admin\AppData\Local\Temp\tmp*.pptx")
    
    if temp_files:
        latest = max(temp_files, key=os.path.getmtime)
        print(f"📁 Analizando: {latest}")
        analyze_logo(latest)
    else:
        print("❌ No se encontró archivo PPTX en temp")
        print("💡 Sube un template en la app primero")

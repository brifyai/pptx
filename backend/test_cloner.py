"""Test del clonador XML"""
import sys
import os

# Agregar paths necesarios
backend_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(backend_dir)
if parent_dir not in sys.path:
    sys.path.insert(0, parent_dir)

# Test 1: Imports
print("=" * 50)
print("TEST 1: Imports")
print("=" * 50)
try:
    from pptx_generator import XML_CLONER_AVAILABLE, generate_presentation
    print(f"XML_CLONER_AVAILABLE: {XML_CLONER_AVAILABLE}")
    assert XML_CLONER_AVAILABLE == True, "Clonador no disponible"
    print("PASS")
except Exception as e:
    print(f"FAIL: {e}")
    sys.exit(1)

# Test 2: Importar pptx_xml_cloner
print("\n" + "=" * 50)
print("TEST 2: Importar pptx_xml_cloner")
print("=" * 50)
try:
    from pptx_xml_cloner import clone_pptx_preserving_all, PPTXXMLCloner
    print("Clonador importado OK")
    print("PASS")
except Exception as e:
    print(f"FAIL: {e}")
    sys.exit(1)

# Test 3: Generar presentación
print("\n" + "=" * 50)
print("TEST 3: Generación de presentación")
print("=" * 50)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import tempfile
    
    # Crear template simple
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    title_box.text_frame.paragraphs[0].text = "Click to add title"
    title_box.text_frame.paragraphs[0].font.size = Pt(44)
    title_box.text_frame.paragraphs[0].font.bold = True
    
    template_path = tempfile.mktemp(suffix='.pptx')
    prs.save(template_path)
    print(f"Template creado: {template_path}")
    
    # Generar con IA
    test_content = {
        'slides': [
            {'content': {'title': 'Titulo Generado', 'subtitle': 'Subtitulo'}},
            {'content': {'title': 'Slide 2', 'heading': 'Seccion 2', 'bullets': ['Punto 1', 'Punto 2']}}
        ]
    }
    
    output = generate_presentation(template_path, test_content)
    print(f"Output: {output}")
    print(f"Existe: {os.path.exists(output)}")
    
    if os.path.exists(output):
        size_kb = os.path.getsize(output) / 1024
        print(f"Tamano: {size_kb:.1f} KB")
        os.unlink(output)
        print("PASS")
    else:
        print("FAIL: Output no existe")
        sys.exit(1)
    
    os.unlink(template_path)
    
except Exception as e:
    import traceback
    print(f"FAIL: {e}")
    traceback.print_exc()
    sys.exit(1)

# Test 4: Verificar preservación de VBA
print("\n" + "=" * 50)
print("TEST 4: PPTXXMLCloner con VBA")
print("=" * 50)
try:
    from pptx_xml_cloner import PPTXXMLCloner
    import tempfile
    from pptx import Presentation
    from pptx.util import Inches
    
    # Crear template
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    title_box.text_frame.paragraphs[0].text = "Click to add title"
    
    template_path = tempfile.mktemp(suffix='.pptx')
    prs.save(template_path)
    
    # Analizar template
    cloner = PPTXXMLCloner(template_path)
    print(f"Slides detectados: {cloner.slide_count}")
    print(f"Has VBA macros: {cloner.has_vba_macros}")
    print(f"Info: {cloner.get_template_info()}")
    
    os.unlink(template_path)
    print("PASS")
    
except Exception as e:
    import traceback
    print(f"FAIL: {e}")
    traceback.print_exc()
    sys.exit(1)

print("\n" + "=" * 50)
print("TODOS LOS TESTS PASARON")
print("=" * 50)
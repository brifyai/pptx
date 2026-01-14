"""
Preservador de tablas para PPTX
Permite preservar y modificar tablas usando manipulación XML directa
"""
from pptx import Presentation
from pptx.table import Table, _Cell
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.oxml.table import CT_Table, CT_TableRow, CT_TableCell
from lxml import etree
from typing import Dict, List, Any, Optional, Tuple
import re


def get_table_xml(table: Table) -> str:
    """
    Obtiene el XML de una tabla.
    
    Args:
        table: Objeto Table de python-pptx
    
    Returns:
        String con el XML de la tabla
    """
    try:
        if hasattr(table, '_tbl'):
            return etree.tostring(table._tbl, encoding='unicode')
        elif hasattr(table, '_element'):
            return etree.tostring(table._element, encoding='unicode')
    except Exception as e:
        print(f"⚠️ Error obteniendo XML de tabla: {e}")
    return ""


def parse_table_xml(xml_str: str) -> Optional[etree._Element]:
    """
    Parsea el XML de una tabla.
    
    Args:
        xml_str: String con XML de tabla
    
    Returns:
        Elemento XML parseado
    """
    try:
        return etree.fromstring(xml_str)
    except Exception as e:
        print(f"⚠️ Error parseando XML de tabla: {e}")
    return None


def extract_table_data(table: Table) -> Dict[str, Any]:
    """
    Extrae todos los datos de una tabla.
    
    Args:
        table: Objeto Table de python-pptx
    
    Returns:
        Diccionario con datos de la tabla
    """
    data = {
        'rows': 0,
        'cols': 0,
        'cells': [],
        'merged_cells': [],
        'styles': {}
    }
    
    try:
        # Dimensiones
        data['rows'] = len(table.rows)
        data['cols'] = len(table.columns)
        
        # Extraer celdas
        for row_idx, row in enumerate(table.rows):
            row_data = []
            for col_idx, cell in enumerate(row.cells):
                cell_data = {
                    'text': '',
                    'row': row_idx,
                    'col': col_idx,
                    'merge_down': 0,
                    'merge_across': 0,
                    'style': {}
                }
                
                # Texto de la celda
                if hasattr(cell, 'text_frame') and cell.text_frame:
                    cell_data['text'] = cell.text_frame.text
                
                # Verificar si es parte de celdas fusionadas
                if hasattr(cell, '_tc'):
                    tc = cell._tc
                    
                    # Verificar gridSpan (fusion horizontal)
                    if hasattr(tc, 'gridSpan'):
                        span = tc.gridSpan
                        if hasattr(span, 'val'):
                            cell_data['merge_across'] = int(span.val) - 1
                    
                    # Verificar vMerge (fusion vertical)
                    if hasattr(tc, 'vMerge'):
                        vmerge = tc.vMerge
                        if hasattr(vmerge, 'val'):
                            cell_data['merge_down'] = int(vmerge.val) if vmerge.val else 0
                
                # Estilo
                if hasattr(cell, 'fill') and cell.fill:
                    fill = cell.fill
                    if hasattr(fill, 'fore_color') and fill.fore_color:
                        rgb = fill.fore_color.rgb
                        if rgb:
                            cell_data['style']['fill_color'] = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                
                if hasattr(cell, 'text_frame') and cell.text_frame:
                    tf = cell.text_frame
                    if hasattr(tf, 'paragraphs') and tf.paragraphs:
                        para = tf.paragraphs[0]
                        if hasattr(para, 'font') and para.font:
                            font = para.font
                            if hasattr(font, 'size') and font.size:
                                cell_data['style']['font_size'] = font.size.pt
                            if hasattr(font, 'bold') and font.bold:
                                cell_data['style']['bold'] = font.bold
                            if hasattr(font, 'color') and font.color:
                                rgb = font.color.rgb
                                if rgb:
                                    cell_data['style']['text_color'] = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                
                row_data.append(cell_data)
            
            data['cells'].append(row_data)
        
        # Detectar celdas fusionadas
        for r_idx, row in enumerate(data['cells']):
            for c_idx, cell in enumerate(row):
                if cell['merge_across'] > 0 or cell['merge_down'] > 0:
                    data['merged_cells'].append({
                        'row': r_idx,
                        'col': c_idx,
                        'merge_down': cell['merge_down'],
                        'merge_across': cell['merge_across']
                    })
        
    except Exception as e:
        print(f"⚠️ Error extrayendo datos de tabla: {e}")
    
    return data


def generate_table_xml(table_data: Dict[str, Any]) -> str:
    """
    Genera XML de tabla a partir de datos.
    
    Args:
        table_data: Datos de la tabla
    
    Returns:
        String con XML de tabla
    """
    rows = table_data.get('rows', 0)
    cols = table_data.get('cols', 0)
    cells = table_data.get('cells', [])
    
    # Namespace de PPTX
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }
    
    # Crear elemento tabla
    table_xml = f'''
    <a:tbl xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:tblPr>
            <a:tblW w="0" type="auto"/>
            <a:tblLook firstRow="1" lastRow="0" firstColumn="1" lastColumn="0" noHBand="0" noVBand="0"/>
        </a:tblPr>
    '''
    
    # Agregar filas y celdas
    for r_idx in range(rows):
        table_xml += '        <a:tr height="0">\n'
        
        for c_idx in range(cols):
            cell_text = ""
            cell_style = {}
            
            if r_idx < len(cells) and c_idx < len(cells[r_idx]):
                cell_text = cells[r_idx][c_idx].get('text', '')
                cell_style = cells[r_idx][c_idx].get('style', {})
            
            # Construir celda
            cell_xml = f'''
            <a:tc>
                <a:txPr>
                    <a:bodyPr/>
                    <a:p>
                        <a:r>
                            <a:rPr'''

            
            # Estilo del texto
            if cell_style.get('bold'):
                cell_xml += ' b="1"'
            
            if cell_style.get('text_color'):
                color = cell_style['text_color'].lstrip('#')
                cell_xml += f' sz="{cell_style.get("font_size", 12) * 100}"'
                cell_xml += f'><a:srgbClr val="{color}"'
                if cell_style.get('text_color'):
                    cell_xml += f'><a:srgbClr val="{color}"/></a:srgbClr>'
                cell_xml += '</a:srgbClr></a:srgbClr>'
            else:
                cell_xml += f' sz="{cell_style.get("font_size", 12) * 100}"/>'
            
            cell_xml += f'''
                            <a:t>{cell_text}</a:t>
                        </a:r>
                    </a:p>
                </a:txPr>
            '''
            
            # Color de fondo
            if cell_style.get('fill_color'):
                color = cell_style['fill_color'].lstrip('#')
                cell_xml += f'<a:tcPr><a:fill><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:fill></a:tcPr>\n'
            
            cell_xml += '            </a:tc>\n'
        
        table_xml += '        </a:tr>\n'
    
    table_xml += '    </a:tbl>'
    
    return table_xml


def update_table_with_data(table: Table, table_data: Dict[str, Any]) -> bool:
    """
    Actualiza una tabla con nuevos datos.
    
    Args:
        table: Objeto Table de python-pptx
        table_data: Nuevos datos para la tabla
    
    Returns:
        True si se actualizó correctamente
    """
    try:
        cells = table_data.get('cells', [])
        
        for r_idx, row in enumerate(table.rows):
            if r_idx >= len(cells):
                break
            
            for c_idx, cell in enumerate(row.cells):
                if c_idx >= len(cells[r_idx]):
                    break
                
                cell_data = cells[r_idx][c_idx]
                new_text = cell_data.get('text', '')
                
                # Actualizar texto
                if hasattr(cell, 'text_frame') and cell.text_frame:
                    cell.text_frame.clear()
                    if new_text:
                        para = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else cell.text_frame.add_paragraph()
                        para.text = new_text
                
                # Actualizar estilo
                style = cell_data.get('style', {})
                if style:
                    if hasattr(cell, 'fill') and cell.fill:
                        if 'fill_color' in style:
                            # El color viene como #RRGGBB
                            color_hex = style['fill_color'].lstrip('#')
                            try:
                                from pptx.dml.color import RGBColor
                                fill = cell.fill
                                fill.solid()
                                fill.fore_color.rgb = RGBColor(
                                    int(color_hex[0:2], 16),
                                    int(color_hex[2:4], 16),
                                    int(color_hex[4:6], 16)
                                )
                            except:
                                pass
        
        return True
        
    except Exception as e:
        print(f"⚠️ Error actualizando tabla: {e}")
        return False


def create_table_from_data(prs, slide, left, top, width, height, 
                          table_data: Dict[str, Any]) -> Optional[Table]:
    """
    Crea una nueva tabla con los datos proporcionados.
    
    Args:
        prs: Presentación
        slide: Slide donde agregar la tabla
        left, top, width, height: Posición y tamaño
        table_data: Datos de la tabla
    
    Returns:
        Objeto Table creado
    """
    try:
        rows = table_data.get('rows', 3)
        cols = table_data.get('cols', 3)
        cells = table_data.get('cells', [])
        
        # Crear tabla
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Llenar datos
        for r_idx in range(rows):
            for c_idx in range(cols):
                cell = table.cell(r_idx, c_idx)
                
                if r_idx < len(cells) and c_idx < len(cells[r_idx]):
                    cell_data = cells[r_idx][c_idx]
                    text = cell_data.get('text', '')
                    style = cell_data.get('style', {})
                    
                    # Texto
                    if text:
                        cell.text = text
                    
                    # Estilo
                    if style:
                        try:
                            if 'fill_color' in style:
                                color_hex = style['fill_color'].lstrip('#')
                                from pptx.dml.color import RGBColor
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(
                                    int(color_hex[0:2], 16),
                                    int(color_hex[2:4], 16),
                                    int(color_hex[4:6], 16)
                                )
                        except:
                            pass
        
        return table
        
    except Exception as e:
        print(f"⚠️ Error creando tabla: {e}")
        return None


def analyze_table_for_ai(table_data: Dict[str, Any]) -> str:
    """
    Genera una descripción de la tabla para enviar a la IA.
    
    Args:
        table_data: Datos de la tabla
    
    Returns:
        Descripción textual de la tabla
    """
    description = []
    
    description.append(f"Tabla de {table_data.get('rows', 0)} filas x {table_data.get('cols', 0)} columnas")
    
    # Celdas fusionadas
    merged = table_data.get('merged_cells', [])
    if merged:
        description.append(f"Celdas fusionadas: {len(merged)}")
        for m in merged[:3]:
            description.append(f"   - Fila {m['row']}, Columna {m['col']}: "
                             f"{m['merge_across']} horizontal, {m['merge_down']} vertical")
    
    # Contenido
    cells = table_data.get('cells', [])
    if cells:
        description.append("\nContenido:")
        for r_idx, row in enumerate(cells[:5]):
            row_text = []
            for c_idx, cell in enumerate(row[:5]):
                text = cell.get('text', '')[:20]
                row_text.append(text if text else '·')
            description.append(f"   Fila {r_idx}: {' | '.join(row_text)}")
        
        if len(cells) > 5:
            description.append(f"   ... y {len(cells) - 5} filas más")
    
    return '\n'.join(description)


def preserve_table_xml(table: Table) -> Dict[str, Any]:
    """
    Preserva el XML completo de una tabla para mantener todas las propiedades.
    
    Args:
        table: Objeto Table de python-pptx
    
    Returns:
        Diccionario con XML preservado
    """
    preservation = {
        'xml': '',
        'properties': {}
    }
    
    try:
        # Obtener XML original
        if hasattr(table, '_tbl'):
            preservation['xml'] = etree.tostring(table._tbl, encoding='unicode', pretty_print=True)
        
        # Propiedades de estilo
        if hasattr(table, '_tbl'):
            tbl = table._tbl
            
            # tblPr (propiedades de tabla)
            tbl_pr = tbl.find('.//a:tblPr', namespaces=nsmap)
            if tbl_pr is not None:
                preservation['properties']['tblPr'] = etree.tostring(tbl_pr, encoding='unicode')
            
            # tblGrid (definición de columnas)
            tbl_grid = tbl.find('.//a:tblGrid', namespaces=nsmap)
            if tbl_grid is not None:
                preservation['properties']['tblGrid'] = etree.tostring(tbl_grid, encoding='unicode')
        
    except Exception as e:
        print(f"⚠️ Error preservando XML de tabla: {e}")
    
    return preservation


def restore_table_from_preservation(table: Table, preservation: Dict[str, Any]) -> bool:
    """
    Restaura el XML preservado en una tabla.
    
    Args:
        table: Objeto Table de python-pptx
        preservation: Datos preservados
    
    Returns:
        True si se restauró correctamente
    """
    try:
        xml = preservation.get('xml', '')
        if not xml:
            return False
        
        # Parsear XML
        new_element = parse_xml(xml)
        
        # Reemplazar elemento
        if hasattr(table, '_element') and hasattr(table._element, 'getparent'):
            parent = table._element.getparent()
            if parent is not None:
                index = parent.index(table._element)
                parent.remove(table._element)
                parent.insert(index, new_element)
                return True
        
        # Intentar reemplazar _tbl
        if hasattr(table, '_tbl'):
            table._tbl.getparent().replace(table._tbl, new_element)
            return True
        
    except Exception as e:
        print(f"⚠️ Error restaurando tabla: {e}")
    
    return False


# Función de prueba
if __name__ == '__main__':
    import sys
    
    if len(sys.argv) < 2:
        print("Uso: python table_preserver.py <archivo.pptx>")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    
    print(f"📋 Analizando tablas en: {pptx_path}")
    
    prs = Presentation(pptx_path)
    table_count = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'has_table') and shape.has_table:
                table_count += 1
                table = shape.table
                data = extract_table_data(table)
                
                print(f"\n--- Tabla {table_count} (Slide {slide_idx + 1}) ---")
                print(f"Dimensiones: {data['rows']} x {data['cols']}")
                print(f"Celdas fusionadas: {len(data['merged_cells'])}")
                
                # Mostrar contenido
                cells = data['cells']
                for r_idx, row in enumerate(cells[:3]):
                    row_text = []
                    for c_idx, cell in enumerate(row[:4]):
                        text = cell.get('text', '')[:15]
                        row_text.append(text if text else '·')
                    print(f"   {' | '.join(row_text)}")
    
    print(f"\n📋 Total de tablas encontradas: {table_count}")
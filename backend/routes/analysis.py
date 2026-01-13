"""
Analysis routes - PPTX/PDF analysis, fonts, content extraction.
"""
from fastapi import APIRouter, UploadFile, File, HTTPException
import tempfile
import os
import json

from pptx import Presentation
from pptx_analyzer import analyze_presentation
from font_detector import analyze_fonts as analyze_pptx_fonts
from utils.logging_utils import logger

router = APIRouter(prefix="/api", tags=["analysis"])


def analyze_pdf(pdf_path: str) -> dict:
    """
    Analiza un archivo PDF y extrae sus páginas como slides.
    Intenta usar PyMuPDF (fitz) primero, luego pdf2image como fallback.
    """
    import base64
    from io import BytesIO
    
    # Intentar con PyMuPDF (no requiere Poppler)
    try:
        import fitz  # PyMuPDF
        
        logger.info(f"📄 Convirtiendo PDF con PyMuPDF: {pdf_path}")
        
        doc = fitz.open(pdf_path)
        slides = []
        slide_images = []
        
        for i, page in enumerate(doc):
            # Renderizar página a imagen (zoom 2x para mejor calidad)
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            
            # Convertir a base64
            img_data = pix.tobytes("png")
            img_base64 = base64.b64encode(img_data).decode()
            
            slide_images.append(f"data:image/png;base64,{img_base64}")
            
            slides.append({
                "number": i + 1,
                "type": "content",
                "layout": "PDF Page",
                "layoutType": "other",
                "isTitle": i == 0,
                "isCover": i == 0,
                "textAreas": [],
                "preview": f"data:image/png;base64,{img_base64}"
            })
        
        # Obtener tamaño de la primera página
        if doc.page_count > 0:
            page = doc[0]
            width = int(page.rect.width * 2)  # Con zoom 2x
            height = int(page.rect.height * 2)
        else:
            width, height = 1920, 1080
        
        doc.close()
        
        return {
            "fileName": os.path.basename(pdf_path),
            "slideSize": {
                "width": width * 9525,
                "height": height * 9525
            },
            "slides": slides,
            "slideImages": slide_images,
            "extractedAssets": {
                "logos": [],
                "images": [],
                "transparentImages": [],
                "animatedElements": []
            },
            "sourceType": "pdf"
        }
        
    except ImportError:
        logger.info("PyMuPDF no disponible, intentando con pdf2image...")
    
    # Fallback a pdf2image (requiere Poppler)
    try:
        from pdf2image import convert_from_path
        from PIL import Image
        
        logger.info(f"📄 Convirtiendo PDF con pdf2image: {pdf_path}")
        
        # Convertir PDF a imágenes
        images = convert_from_path(pdf_path, dpi=150)
        
        slides = []
        slide_images = []
        
        for i, image in enumerate(images):
            # Convertir imagen a base64
            buffered = BytesIO()
            image.save(buffered, format="PNG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            
            slide_images.append(f"data:image/png;base64,{img_base64}")
            
            slides.append({
                "number": i + 1,
                "type": "content",
                "layout": "PDF Page",
                "layoutType": "other",
                "isTitle": i == 0,
                "isCover": i == 0,
                "textAreas": [],
                "preview": f"data:image/png;base64,{img_base64}"
            })
        
        width, height = images[0].size if images else (1920, 1080)
        
        return {
            "fileName": os.path.basename(pdf_path),
            "slideSize": {
                "width": width * 9525,
                "height": height * 9525
            },
            "slides": slides,
            "slideImages": slide_images,
            "extractedAssets": {
                "logos": [],
                "images": [],
                "transparentImages": [],
                "animatedElements": []
            },
            "sourceType": "pdf"
        }
        
    except ImportError:
        logger.warning("⚠️ Ni PyMuPDF ni pdf2image están instalados.")
    except Exception as e:
        logger.warning(f"⚠️ Error con pdf2image: {e}")
    
    # Fallback final: placeholder
    logger.warning("⚠️ Usando análisis básico de PDF (sin imágenes).")
    return {
        "fileName": os.path.basename(pdf_path),
        "slideSize": {"width": 9144000, "height": 6858000},
        "slides": [{
            "number": 1,
            "type": "content",
            "layout": "PDF",
            "layoutType": "other",
            "isTitle": True,
            "isCover": True,
            "textAreas": [],
            "preview": None
        }],
        "slideImages": [],
        "extractedAssets": {
            "logos": [],
            "images": [],
            "transparentImages": [],
            "animatedElements": []
        },
        "sourceType": "pdf",
        "warning": "Instala PyMuPDF (pip install pymupdf) para soporte completo de PDF"
    }


@router.post("/analyze")
async def analyze_ppt(file: UploadFile = File(...)):
    """
    Analiza un archivo PowerPoint o PDF y extrae toda su estructura de diseño.
    """
    filename = file.filename.lower()
    
    # Determinar tipo de archivo
    if filename.endswith('.pdf'):
        file_type = 'pdf'
        suffix = '.pdf'
    elif filename.endswith(('.pptx', '.ppt')):
        file_type = 'pptx'
        suffix = '.pptx'
    else:
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx o .pdf")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name
        
        logger.info(f"📄 Archivo guardado en: {tmp_path} (tipo: {file_type})")
        
        if file_type == 'pdf':
            analysis = analyze_pdf(tmp_path)
        else:
            analysis = analyze_presentation(tmp_path)
        
        logger.info(f"✅ Análisis completado: {len(analysis['slides'])} slides")
        
        os.unlink(tmp_path)
        
        return {
            "success": True,
            "analysis": analysis,
            "message": f"Análisis completado: {len(analysis['slides'])} diapositivas detectadas"
        }
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al analizar: {str(e)}")


@router.post("/analyze-fonts")
async def analyze_fonts_endpoint(file: UploadFile = File(...)):
    """
    Analiza las fuentes usadas en un PPTX y detecta cuáles faltan.
    """
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name
        
        logger.info(f"🔤 Analizando fuentes de: {file.filename}")
        
        font_analysis = analyze_pptx_fonts(tmp_path)
        
        os.unlink(tmp_path)
        
        logger.info(f"✅ Análisis de fuentes completado: {font_analysis['summary']}")
        
        return {
            "success": True,
            "fileName": file.filename,
            **font_analysis
        }
    
    except Exception as e:
        logger.error(f"❌ Error analizando fuentes: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al analizar fuentes: {str(e)}")


@router.post("/extract-content")
async def extract_content(file: UploadFile = File(...)):
    """
    Extrae solo el contenido de texto de un PPTX (sin diseño).
    """
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = tmp.name
        
        logger.info(f"📄 Extrayendo contenido de: {tmp_path}")
        
        prs = Presentation(tmp_path)
        
        extracted_slides = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = {
                "slideNumber": slide_idx + 1,
                "type": "content",
                "texts": []
            }
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        text_type = "body"
                        if shape.is_placeholder:
                            placeholder_type = shape.placeholder_format.type
                            if placeholder_type == 1:
                                text_type = "title"
                            elif placeholder_type == 2:
                                text_type = "subtitle"
                        
                        if '\n' in text and text_type == "body":
                            bullets = [line.strip() for line in text.split('\n') if line.strip()]
                            slide_content["texts"].append({
                                "type": "bullets",
                                "content": bullets
                            })
                        else:
                            slide_content["texts"].append({
                                "type": text_type,
                                "content": text
                            })
            
            if slide_idx == 0 or any(t["type"] == "title" for t in slide_content["texts"]):
                slide_content["type"] = "title"
            
            extracted_slides.append(slide_content)
        
        os.unlink(tmp_path)
        
        logger.info(f"✅ Contenido extraído: {len(extracted_slides)} slides")
        
        return {
            "success": True,
            "fileName": file.filename,
            "slideCount": len(extracted_slides),
            "slides": extracted_slides
        }
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al extraer contenido: {str(e)}")

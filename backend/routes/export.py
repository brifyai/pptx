"""
Export routes - PPTX and PDF export functionality.
Supports both sync (small files) and async (large files) generation.
"""
from fastapi import APIRouter, UploadFile, File, HTTPException, Request, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse
import tempfile
import os
import json
import re
import base64
import asyncio

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx_generator import generate_presentation
from utils.logging_utils import logger
from core.task_queue import task_queue, TaskStatus

router = APIRouter(prefix="/api", tags=["export"])


# ============================================
# ASYNC TASK ENDPOINTS (for heavy operations)
# ============================================

@router.post("/export/pptx/async")
async def export_pptx_async(request: Request, background_tasks: BackgroundTasks):
    """
    Async PPTX export - returns task ID immediately.
    Use GET /api/task/{task_id} to poll status.
    Use GET /api/task/{task_id}/download when completed.
    
    For large files or high concurrency scenarios.
    """
    try:
        # Parse request (same as sync version)
        content_type = request.headers.get('content-type', '')
        template_content = None
        template_filename = None
        data = None
        
        if 'multipart/form-data' in content_type:
            body = await request.body()
            boundary_match = re.search(r'boundary=([^;]+)', content_type)
            if not boundary_match:
                raise HTTPException(status_code=400, detail="Invalid multipart boundary")
            
            boundary = boundary_match.group(1).strip('"')
            parts = body.split(f'--{boundary}'.encode())
            
            for part in parts:
                if b'name="template"' in part:
                    filename_match = re.search(rb'filename="([^"]+)"', part)
                    if filename_match:
                        template_filename = filename_match.group(1).decode('utf-8')
                    content_start = part.find(b'\r\n\r\n')
                    if content_start != -1:
                        template_content = part[content_start + 4:]
                        if template_content.endswith(b'\r\n'):
                            template_content = template_content[:-2]
                elif b'name="data"' in part:
                    content_start = part.find(b'\r\n\r\n')
                    if content_start != -1:
                        data_bytes = part[content_start + 4:]
                        if data_bytes.endswith(b'\r\n'):
                            data_bytes = data_bytes[:-2]
                        data = data_bytes.decode('utf-8')
        else:
            body = await request.body()
            json_data = json.loads(body)
            data = json.dumps(json_data)
        
        # Create task
        task = task_queue.create_task()
        
        # Save template to temp file for background processing
        template_path = None
        if template_content:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                tmp.write(template_content)
                template_path = tmp.name
        
        # Schedule background task
        asyncio.create_task(
            task_queue.run_in_background(
                task.id,
                _generate_pptx_task,
                template_path,
                data
            )
        )
        
        logger.info(f"📋 PPTX generation queued: task {task.id}")
        
        return JSONResponse({
            "success": True,
            "taskId": task.id,
            "status": "pending",
            "message": "Generación en cola. Usa GET /api/task/{taskId} para verificar estado."
        })
        
    except Exception as e:
        logger.error(f"❌ Error queuing PPTX export: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/task/{task_id}")
async def get_task_status(task_id: str):
    """Get status of an async task."""
    task = task_queue.get_task(task_id)
    
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    response = task.to_dict()
    
    # If completed, include download URL
    if task.status == TaskStatus.COMPLETED and task.result:
        response["downloadUrl"] = f"/api/task/{task_id}/download"
    
    return response


@router.get("/task/{task_id}/download")
async def download_task_result(task_id: str):
    """Download the result of a completed task."""
    task = task_queue.get_task(task_id)
    
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    if task.status != TaskStatus.COMPLETED:
        raise HTTPException(status_code=400, detail=f"Task not completed. Status: {task.status.value}")
    
    if not task.result or not os.path.exists(task.result):
        raise HTTPException(status_code=404, detail="Result file not found")
    
    return FileResponse(
        task.result,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="presentacion.pptx",
        headers={"Content-Disposition": "attachment; filename=presentacion.pptx"}
    )


@router.get("/queue/status")
async def get_queue_status():
    """Get current task queue statistics."""
    return task_queue.get_queue_status()


def _generate_pptx_task(template_path: str, data: str) -> str:
    """
    Synchronous PPTX generation task (runs in thread pool).
    Returns path to generated file.
    """
    try:
        slides = []
        if data:
            slides_data = json.loads(data) if isinstance(data, str) else data
            slides = slides_data.get('slides', [])
        
        if template_path:
            ai_content = {
                'slides': [{'content': slide.get('content', {})} for slide in slides]
            }
            output_path = generate_presentation(template_path, ai_content)
            
            # Cleanup template
            if os.path.exists(template_path):
                os.unlink(template_path)
            
            return output_path
        else:
            # Generate basic presentation without template
            return _generate_basic_pptx(slides)
            
    except Exception as e:
        logger.error(f"❌ PPTX generation failed: {e}")
        raise


def _generate_basic_pptx(slides: list) -> str:
    """Generate a basic PPTX without template."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        slide_type = slide_data.get('type', 'content')
        content = slide_data.get('content', {})
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        
        if slide_type == 'title':
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = content.get('title', '')
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.alignment = 1
        else:
            if content.get('heading'):
                head_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
                head_frame = head_box.text_frame
                head_para = head_frame.paragraphs[0]
                head_para.text = content.get('heading', '')
                head_para.font.size = Pt(32)
                head_para.font.bold = True
    
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        prs.save(tmp.name)
        return tmp.name


# ============================================
# SYNC ENDPOINTS (for small/quick operations)
# ============================================


@router.post("/generate")
async def generate_ppt(file: UploadFile = File(...), content: str = None):
    """
    Genera un nuevo PowerPoint manteniendo el diseño original
    pero con contenido generado por IA.
    """
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            file_content = await file.read()
            tmp.write(file_content)
            original_path = tmp.name
        
        ai_content = json.loads(content) if content else None
        
        output_path = generate_presentation(original_path, ai_content)
        
        os.unlink(original_path)
        
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="presentacion_generada.pptx",
            headers={"Content-Disposition": "attachment; filename=presentacion_generada.pptx"}
        )
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al generar: {str(e)}")


@router.post("/export/pptx")
async def export_pptx(request: Request):
    """
    Exporta slides a PowerPoint (.pptx).
    Si se proporciona un template, clona su diseño y solo reemplaza el texto.
    """
    try:
        content_type = request.headers.get('content-type', '')
        
        template_content = None
        template_filename = None
        data = None
        
        if 'multipart/form-data' in content_type:
            body = await request.body()
            
            boundary_match = re.search(r'boundary=([^;]+)', content_type)
            if not boundary_match:
                raise HTTPException(status_code=400, detail="Invalid multipart boundary")
            
            boundary = boundary_match.group(1).strip('"')
            parts = body.split(f'--{boundary}'.encode())
            
            for part in parts:
                if b'name="template"' in part:
                    filename_match = re.search(rb'filename="([^"]+)"', part)
                    if filename_match:
                        template_filename = filename_match.group(1).decode('utf-8')
                    
                    content_start = part.find(b'\r\n\r\n')
                    if content_start != -1:
                        template_content = part[content_start + 4:]
                        if template_content.endswith(b'\r\n'):
                            template_content = template_content[:-2]
                
                elif b'name="data"' in part:
                    content_start = part.find(b'\r\n\r\n')
                    if content_start != -1:
                        data_bytes = part[content_start + 4:]
                        if data_bytes.endswith(b'\r\n'):
                            data_bytes = data_bytes[:-2]
                        data = data_bytes.decode('utf-8')
            
            logger.info(f"📤 Export PPTX - Template: {template_filename}")
            logger.info(f"📤 Export PPTX - Template size: {len(template_content) if template_content else 0} bytes")
        else:
            body = await request.body()
            json_data = json.loads(body)
            data = json.dumps(json_data)
        
        slides = []
        if data:
            slides_data = json.loads(data) if isinstance(data, str) else data
            slides = slides_data.get('slides', [])
            logger.info(f"📤 Slides parseados: {len(slides)}")
            
            # DEBUG: Mostrar contenido de cada slide
            for idx, slide in enumerate(slides):
                content = slide.get('content', {})
                logger.info(f"📤 Slide {idx+1} contenido:")
                logger.info(f"   - title: {content.get('title', 'N/A')}")
                logger.info(f"   - subtitle: {content.get('subtitle', 'N/A')}")
                logger.info(f"   - heading: {content.get('heading', 'N/A')}")
                logger.info(f"   - bullets: {len(content.get('bullets', []))} items")
                if content.get('bullets'):
                    for i, bullet in enumerate(content.get('bullets', [])[:3]):
                        logger.info(f"     • {bullet}")
        
        # Si hay template, usar clonación completa
        if template_content and template_filename:
            logger.info(f"📄 Usando template: {template_filename}")
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                tmp.write(template_content)
                template_path = tmp.name
            
            ai_content = {
                'slides': [{'content': slide.get('content', {})} for slide in slides]
            }
            
            output_path = generate_presentation(template_path, ai_content)
            
            logger.info(f"✅ Presentación generada: {output_path}")
            
            os.unlink(template_path)
            
            return FileResponse(
                output_path,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                filename="presentacion.pptx",
                headers={"Content-Disposition": "attachment; filename=presentacion.pptx"}
            )
        
        # Si no hay template, crear presentación básica
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for slide_data in slides:
            slide_type = slide_data.get('type', 'content')
            content = slide_data.get('content', {})
            
            layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(layout)
            
            if slide_type == 'title':
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = content.get('title', '')
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.alignment = 1
                
                if content.get('subtitle'):
                    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
                    sub_frame = sub_box.text_frame
                    sub_para = sub_frame.paragraphs[0]
                    sub_para.text = content.get('subtitle', '')
                    sub_para.font.size = Pt(24)
                    sub_para.alignment = 1
            else:
                if content.get('heading'):
                    head_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
                    head_frame = head_box.text_frame
                    head_para = head_frame.paragraphs[0]
                    head_para.text = content.get('heading', '')
                    head_para.font.size = Pt(32)
                    head_para.font.bold = True
                
                bullets = content.get('bullets', [])
                if bullets:
                    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5))
                    bullet_frame = bullet_box.text_frame
                    bullet_frame.word_wrap = True
                    
                    for i, bullet in enumerate(bullets):
                        if bullet and bullet.strip():
                            para = bullet_frame.paragraphs[0] if i == 0 else bullet_frame.add_paragraph()
                            para.text = f"• {bullet}"
                            para.font.size = Pt(18)
                            para.space_after = Pt(12)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            prs.save(tmp.name)
            tmp_path = tmp.name
        
        return FileResponse(
            tmp_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="presentacion.pptx",
            headers={"Content-Disposition": "attachment; filename=presentacion.pptx"}
        )
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al exportar PPTX: {str(e)}")


@router.post("/export/pdf")
async def export_pdf(request: Request):
    """
    Exporta slides a PDF usando las imágenes de preview.
    Tamaño exacto 16:9 sin márgenes.
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        from PIL import Image
        from io import BytesIO
        
        PAGE_WIDTH = 960
        PAGE_HEIGHT = 540
        
        content_type = request.headers.get('content-type', '')
        slides = []
        
        if 'multipart/form-data' in content_type:
            body = await request.body()
            
            boundary_match = re.search(r'boundary=([^;]+)', content_type)
            if boundary_match:
                boundary = boundary_match.group(1).strip('"')
                parts = body.split(f'--{boundary}'.encode())
                
                for part in parts:
                    if b'name="data"' in part:
                        content_start = part.find(b'\r\n\r\n')
                        if content_start != -1:
                            data_bytes = part[content_start + 4:]
                            if data_bytes.endswith(b'\r\n'):
                                data_bytes = data_bytes[:-2]
                            data = json.loads(data_bytes.decode('utf-8'))
                            slides = data.get('slides', [])
        else:
            body = await request.body()
            data = json.loads(body)
            slides = data.get('slides', [])
        
        logger.info(f"📄 Exportando PDF 16:9 con {len(slides)} slides")
        
        pdf_path = tempfile.mktemp(suffix='.pdf')
        c = canvas.Canvas(pdf_path, pagesize=(PAGE_WIDTH, PAGE_HEIGHT))
        
        for idx, slide in enumerate(slides):
            preview = slide.get('preview')
            
            if preview and preview.startswith('data:image'):
                try:
                    header, base64_data = preview.split(',', 1)
                    image_data = base64.b64decode(base64_data)
                    
                    img = Image.open(BytesIO(image_data))
                    
                    if img.mode in ('RGBA', 'P'):
                        img = img.convert('RGB')
                    
                    img_buffer = BytesIO()
                    img.save(img_buffer, format='PNG')
                    img_buffer.seek(0)
                    
                    c.drawImage(ImageReader(img_buffer), 0, 0, PAGE_WIDTH, PAGE_HEIGHT)
                    
                    logger.info(f"   ✅ Slide {idx + 1} agregado al PDF")
                    
                except Exception as e:
                    logger.warning(f"   ⚠️ Error procesando imagen slide {idx + 1}: {e}")
                    c.setFont("Helvetica-Bold", 24)
                    c.drawCentredString(PAGE_WIDTH / 2, PAGE_HEIGHT / 2, f"Slide {idx + 1}")
            else:
                c.setFont("Helvetica-Bold", 32)
                content = slide.get('content', {})
                title = content.get('title') or content.get('heading', f'Slide {idx + 1}')
                c.drawCentredString(PAGE_WIDTH / 2, PAGE_HEIGHT - 80, title)
                
                if content.get('subtitle'):
                    c.setFont("Helvetica", 20)
                    c.drawCentredString(PAGE_WIDTH / 2, PAGE_HEIGHT - 120, content['subtitle'])
                
                bullets = content.get('bullets', [])
                if bullets:
                    c.setFont("Helvetica", 14)
                    y_pos = PAGE_HEIGHT - 180
                    for bullet in bullets:
                        if bullet and bullet.strip():
                            c.drawString(80, y_pos, f"• {bullet}")
                            y_pos -= 22
                
                logger.info(f"   📝 Slide {idx + 1} (texto) agregado al PDF")
            
            if idx < len(slides) - 1:
                c.showPage()
        
        c.save()
        
        logger.info(f"✅ PDF 16:9 generado: {pdf_path}")
        
        return FileResponse(
            pdf_path,
            media_type="application/pdf",
            filename="presentacion.pdf",
            headers={"Content-Disposition": "attachment; filename=presentacion.pdf"}
        )
    
    except Exception as e:
        logger.error(f"❌ Error exportando PDF: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error al exportar PDF: {str(e)}")

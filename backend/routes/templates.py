"""
Template routes - Template analysis, caching, and mapping.
"""
from fastapi import APIRouter, UploadFile, File, HTTPException
import tempfile
import os

from pptx import Presentation
from mapping_cache import MappingCache
from shape_matcher import ShapeMatcher
from schemas.requests import UpdateMappingRequest
from services.gemini_vision import analyze_with_retry, validate_element_type, VALID_ELEMENT_TYPES
from services.slide_converter import convert_slide_to_image
from utils.logging_utils import (
    logger, ErrorCategory, 
    log_error_with_context, log_operation_start, log_operation_success
)

router = APIRouter(prefix="/api", tags=["templates"])
mapping_cache = MappingCache()


@router.post("/analyze-template")
async def analyze_template(file: UploadFile = File(...)):
    """
    Analyze a PPTX template and return mapping of elements.
    
    1. Verifies cache by hash
    2. If not cached, converts to image and calls Gemini Vision
    3. Links detections with real shapes using ShapeMatcher
    4. Saves to cache and returns
    """
    operation_context = {
        "filename": file.filename,
        "content_type": file.content_type
    }
    log_operation_start("analyze_template", operation_context)
    
    if not file.filename.endswith(('.pptx', '.ppt')):
        log_error_with_context(
            ValueError("Invalid file type"),
            ErrorCategory.FILE_VALIDATION,
            "validate_file_type",
            {"filename": file.filename, "expected": [".pptx", ".ppt"]},
            include_traceback=False
        )
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx")
    
    tmp_path = None
    
    try:
        file_content = await file.read()
        operation_context["file_size"] = len(file_content)
        
        template_hash = mapping_cache.generate_template_hash(file_content)
        operation_context["template_hash"] = template_hash[:16] + "..."
        logger.info(f"📄 Template hash: {template_hash[:16]}...")
        
        # Check cache first
        cached_mapping = mapping_cache.get_cached_mapping(template_hash)
        if cached_mapping:
            log_operation_success("analyze_template", f"Loaded from cache - {len(cached_mapping['elements'])} elements")
            return {
                "success": True,
                "templateHash": template_hash,
                "elements": cached_mapping['elements'],
                "shapeMapping": cached_mapping['shapeMapping'],
                "source": "cache",
                "message": f"Template reconocido - {len(cached_mapping['elements'])} elementos cargados desde caché"
            }
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name
        
        logger.info(f"📄 Archivo guardado temporalmente: {tmp_path}")
        
        try:
            prs = Presentation(tmp_path)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            operation_context["slide_dimensions"] = f"{slide_width}x{slide_height}"
        except Exception as e:
            log_error_with_context(e, ErrorCategory.FILE_CONVERSION, "open_presentation", operation_context)
            raise HTTPException(status_code=400, detail="No se pudo abrir el archivo PPTX.")
        
        shapes = list(prs.slides[0].shapes) if prs.slides else []
        operation_context["shape_count"] = len(shapes)
        
        logger.info(f"📐 Slide dimensions: {slide_width} x {slide_height} EMUs")
        logger.info(f"🔷 Found {len(shapes)} shapes in first slide")
        
        # Convert slide to image
        logger.info("🎨 Converting slide to image...")
        try:
            slide_image = convert_slide_to_image(tmp_path, 0)
            if not slide_image:
                raise ValueError("Image conversion returned empty result")
        except Exception as e:
            log_error_with_context(e, ErrorCategory.FILE_CONVERSION, "convert_slide_to_image", operation_context)
            raise HTTPException(status_code=500, detail="Error de conversión: No se pudo convertir el slide a imagen.")
        
        # Call Gemini Vision
        logger.info("🔍 Calling Gemini Vision API...")
        try:
            vision_result = await analyze_with_retry(slide_image)
            operation_context["detected_elements"] = len(vision_result.get('elements', []))
        except HTTPException:
            raise
        except Exception as e:
            log_error_with_context(e, ErrorCategory.GEMINI_API, "analyze_with_retry", operation_context)
            raise HTTPException(status_code=500, detail=f"Error de análisis visual: {str(e)}")
        
        # Match detections to shapes
        logger.info("🔗 Matching detections to shapes...")
        try:
            matcher = ShapeMatcher(slide_width, slide_height)
            detections = vision_result['elements']
            shape_mapping = matcher.match_detections_to_shapes(shapes, detections)
            
            for element in detections:
                element_type = element['type']
                if element_type in shape_mapping:
                    element['shapeId'] = shape_mapping[element_type]
            
            for shape in shapes:
                shape_id = getattr(shape, 'shape_id', None)
                if shape_id and shape_id not in shape_mapping.values():
                    classified_type = matcher.classify_shape(shape)
                    if classified_type != 'UNKNOWN' and classified_type not in shape_mapping:
                        shape_mapping[classified_type] = shape_id
            
            operation_context["shape_mapping"] = shape_mapping
            logger.info(f"✅ Shape mapping complete: {shape_mapping}")
            
        except Exception as e:
            log_error_with_context(e, ErrorCategory.SHAPE_MATCHING, "match_detections_to_shapes", operation_context)
            logger.warning("⚠️ Shape matching failed, continuing with empty mapping")
            shape_mapping = {}
        
        # Save to cache
        try:
            mapping_to_save = {'elements': detections, 'shapeMapping': shape_mapping}
            mapping_cache.save_mapping(
                template_hash=template_hash,
                mapping=mapping_to_save,
                template_name=file.filename,
                file_path=None,
                thumbnail_url=slide_image[:100] + '...' if slide_image else None
            )
            logger.info(f"💾 Mapping saved to cache")
        except Exception as e:
            log_error_with_context(e, ErrorCategory.CACHE_OPERATION, "save_mapping", operation_context)
            logger.warning("⚠️ Failed to save mapping to cache, continuing...")
        
        log_operation_success("analyze_template", f"{len(detections)} elements detected")
        
        return {
            "success": True,
            "templateHash": template_hash,
            "elements": detections,
            "shapeMapping": shape_mapping,
            "source": "vision",
            "message": f"Análisis completado - {len(detections)} elementos detectados"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        log_error_with_context(e, ErrorCategory.UNKNOWN, "analyze_template", operation_context)
        raise HTTPException(status_code=500, detail=f"Error al analizar template: {str(e)}")
    
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
                logger.info(f"🗑️ Cleaned up temp file: {tmp_path}")
            except Exception as e:
                logger.warning(f"Could not delete temp file: {e}")


@router.post("/update-mapping")
async def update_mapping(request: UpdateMappingRequest):
    """
    Update element type in cache (manual user correction).
    """
    operation_context = {
        "template_hash": request.template_hash[:16] + "...",
        "element_id": request.element_id,
        "new_type": request.new_type
    }
    log_operation_start("update_mapping", operation_context)
    
    try:
        validated_type = validate_element_type(request.new_type)
        if validated_type == 'UNKNOWN' and request.new_type.upper() != 'UNKNOWN':
            log_error_with_context(
                ValueError(f"Invalid element type: {request.new_type}"),
                ErrorCategory.FILE_VALIDATION,
                "validate_element_type",
                operation_context,
                include_traceback=False
            )
            raise HTTPException(
                status_code=400, 
                detail=f"Invalid element type: {request.new_type}. Valid types: {', '.join(VALID_ELEMENT_TYPES)}"
            )
        
        success = mapping_cache.update_element_type(
            template_hash=request.template_hash,
            element_id=request.element_id,
            new_type=validated_type
        )
        
        if not success:
            log_error_with_context(
                ValueError("Element not found"),
                ErrorCategory.CACHE_OPERATION,
                "update_element_type",
                operation_context,
                include_traceback=False
            )
            raise HTTPException(
                status_code=404, 
                detail=f"Element {request.element_id} not found for template {request.template_hash[:16]}..."
            )
        
        log_operation_success("update_mapping", f"Element {request.element_id} updated to {validated_type}")
        
        return {"success": True, "message": f"Elemento actualizado a tipo {validated_type}"}
        
    except HTTPException:
        raise
    except Exception as e:
        log_error_with_context(e, ErrorCategory.UNKNOWN, "update_mapping", operation_context)
        raise HTTPException(status_code=500, detail=f"Error al actualizar mapping: {str(e)}")


@router.get("/templates")
async def list_templates():
    """List all cached templates."""
    try:
        templates = mapping_cache.get_all_templates()
        return {"success": True, "templates": templates, "count": len(templates)}
    except Exception as e:
        logger.error(f"❌ Error listing templates: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al listar templates: {str(e)}")


@router.get("/template/{template_hash}")
async def get_template_mapping(template_hash: str):
    """Get cached mapping for a specific template by hash."""
    try:
        cached_mapping = mapping_cache.get_cached_mapping(template_hash)
        
        if not cached_mapping:
            raise HTTPException(status_code=404, detail="Template no encontrado en caché")
        
        return {"success": True, "mapping": cached_mapping}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error getting template: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al obtener template: {str(e)}")


@router.delete("/template/{template_hash}")
async def delete_template_mapping(template_hash: str):
    """Delete a cached template mapping."""
    try:
        success = mapping_cache.delete_mapping(template_hash)
        
        if not success:
            raise HTTPException(status_code=404, detail="Template no encontrado en caché")
        
        return {"success": True, "message": "Template eliminado de caché"}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ Error deleting template: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al eliminar template: {str(e)}")

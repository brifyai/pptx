"""
Modificador de gráficos para PPTX
Permite modificar datos de gráficos con contenido generado por IA
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from typing import Dict, List, Any, Optional
import re


def extract_chart_data(chart) -> Dict[str, Any]:
    """
    Extrae todos los datos de un gráfico.
    
    Args:
        chart: Objeto de gráfico de python-pptx
    
    Returns:
        Diccionario con datos del gráfico
    """
    data = {
        'chart_type': None,
        'categories': [],
        'series': [],
        'has_legend': False,
        'legend_position': None,
        'title': None,
        'axis_title_x': None,
        'axis_title_y': None
    }
    
    try:
        # Tipo de gráfico
        if hasattr(chart, 'chart_type'):
            data['chart_type'] = str(chart.chart_type)
        
        # Categorías (eje X)
        if chart.plots and hasattr(chart.plots[0], 'categories') and chart.plots[0].categories:
            data['categories'] = list(chart.plots[0].categories)
        
        # Series de datos
        if hasattr(chart, 'series'):
            for series in chart.series:
                series_data = {
                    'name': series.name if hasattr(series, 'name') else 'Serie',
                    'values': []
                }
                
                # Extraer valores
                if hasattr(series, 'values') and series.values:
                    try:
                        series_data['values'] = [float(v) for v in series.values]
                    except (ValueError, TypeError):
                        series_data['values'] = list(series.values) if series.values else []
                
                data['series'].append(series_data)
        
        # Título del gráfico
        if hasattr(chart, 'chart_title') and chart.chart_title:
            title_frame = chart.chart_title
            if hasattr(title_frame, 'text_frame') and title_frame.text_frame:
                data['title'] = title_frame.text_frame.text
        
        # Leyenda
        if hasattr(chart, 'has_legend'):
            data['has_legend'] = chart.has_legend
        if hasattr(chart, 'legend'):
            legend = chart.legend
            if hasattr(legend, 'position'):
                data['legend_position'] = str(legend.position)
        
        # Títulos de ejes
        if hasattr(chart, 'value_axis') and chart.value_axis:
            axis = chart.value_axis
            if hasattr(axis, 'axis_title') and axis.axis_title:
                if hasattr(axis.axis_title, 'text_frame') and axis.axis_title.text_frame:
                    data['axis_title_y'] = axis.axis_title.text_frame.text
        
        if hasattr(chart, 'category_axis') and chart.category_axis:
            axis = chart.category_axis
            if hasattr(axis, 'axis_title') and axis.axis_title:
                if hasattr(axis.axis_title, 'text_frame') and axis.axis_title.text_frame:
                    data['axis_title_x'] = axis.axis_title.text_frame.text
        
    except Exception as e:
        print(f"⚠️ Error extrayendo datos del gráfico: {e}")
    
    return data


def generate_chart_data_with_ai(chart_data: Dict[str, Any], 
                                content: Dict[str, Any]) -> Dict[str, Any]:
    """
    Genera nuevos datos para un gráfico basándose en contenido IA.
    
    Args:
        chart_data: Datos actuales del gráfico
        content: Contenido generado por IA (puede incluir chart_data)
    
    Returns:
        Nuevos datos para el gráfico
    """
    new_data = chart_data.copy()
    
    # Si hay datos de IA específicos para el gráfico
    if 'chart_data' in content:
        ai_chart_data = content['chart_data']
        
        # Actualizar categorías
        if 'categories' in ai_chart_data:
            new_data['categories'] = ai_chart_data['categories']
        
        # Actualizar series
        if 'series' in ai_chart_data:
            new_data['series'] = ai_chart_data['series']
        
        # Actualizar título
        if 'title' in ai_chart_data:
            new_data['title'] = ai_chart_data['title']
    
    # Generar datos basados en el contenido del slide
    else:
        # Usar bullets como categorías
        bullets = content.get('bullets', [])
        if bullets and not new_data['categories']:
            # Extraer palabras clave de los bullets como categorías
            categories = []
            for bullet in bullets[:5]:
                # Tomar las primeras palabras como categorías
                words = bullet.split()[:2]
                if words:
                    categories.append(' '.join(words))
            if categories:
                new_data['categories'] = categories
        
        # Generar valores basados en el contenido
        if new_data['categories'] and not new_data['series']:
            # Crear una serie con valores basados en la longitud del texto
            values = []
            for cat in new_data['categories']:
                # Valor basado en la relevancia (simulado)
                val = len(cat) * 10
                values.append(val)
            
            new_data['series'] = [{
                'name': 'Datos',
                'values': values
            }]
    
    return new_data


def update_chart_with_data(chart, new_data: Dict[str, Any]) -> bool:
    """
    Actualiza un gráfico con nuevos datos.
    
    Args:
        chart: Objeto de gráfico de python-pptx
        new_data: Nuevos datos para el gráfico
    
    Returns:
        True si se actualizó correctamente
    """
    try:
        # Actualizar categorías
        if 'categories' in new_data and new_data['categories']:
            categories = new_data['categories']
            
            # Crear nuevo ChartData
            new_chart_data = CategoryChartData()
            new_chart_data.categories = categories
            
            # Agregar series
            for series_data in new_data.get('series', []):
                values = series_data.get('values', [])
                # Asegurar que hay suficientes valores
                while len(values) < len(categories):
                    values.append(0)
                new_chart_data.add_series(series_data.get('name', 'Serie'), values[:len(categories)])
            
            # Reemplazar datos del gráfico
            chart.chart_data = new_chart_data
        
        # Actualizar título
        if 'title' in new_data and new_data['title']:
            if hasattr(chart, 'chart_title') and chart.chart_title:
                title_frame = chart.chart_title
                if hasattr(title_frame, 'text_frame'):
                    title_frame.text_frame.text = new_data['title']
        
        return True
        
    except Exception as e:
        print(f"⚠️ Error actualizando gráfico: {e}")
        return False


def create_chart_from_data(prs, slide, left, top, width, height, 
                          chart_type: str, data: Dict[str, Any]):
    """
    Crea un nuevo gráfico con los datos proporcionados.
    
    Args:
        prs: Presentación
        slide: Slide donde agregar el gráfico
        left, top, width, height: Posición y tamaño
        chart_type: Tipo de gráfico
        data: Datos del gráfico
    
    Returns:
        Objeto del gráfico creado
    """
    try:
        # Mapear tipo de string a XL_CHART_TYPE
        chart_type_map = {
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'area': XL_CHART_TYPE.AREA,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
            'doughnut': XL_CHART_TYPE.DOUGHNUT
        }
        
        xl_type = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Crear ChartData
        chart_data = CategoryChartData()
        chart_data.categories = data.get('categories', ['A', 'B', 'C'])
        
        for series in data.get('series', []):
            values = series.get('values', [1, 2, 3])
            chart_data.add_series(series.get('name', 'Serie'), values)
        
        # Agregar gráfico
        chart = slide.shapes.add_chart(
            xl_type,
            left, top, width, height,
            chart_data
        )
        
        # Actualizar título si existe
        if 'title' in data and data['title']:
            try:
                chart.chart_title.text_frame.text = data['title']
            except:
                pass
        
        return chart
        
    except Exception as e:
        print(f"⚠️ Error creando gráfico: {e}")
        return None


def analyze_chart_for_ai(chart_data: Dict[str, Any]) -> str:
    """
    Genera una descripción del gráfico para enviar a la IA.
    
    Args:
        chart_data: Datos del gráfico
    
    Returns:
        Descripción textual del gráfico
    """
    description = []
    
    # Tipo de gráfico
    if chart_data.get('chart_type'):
        description.append(f"Tipo de gráfico: {chart_data['chart_type']}")
    
    # Categorías
    categories = chart_data.get('categories', [])
    if categories:
        description.append(f"Categorías ({len(categories)}): {', '.join(categories[:5])}")
        if len(categories) > 5:
            description.append(f"   ... y {len(categories) - 5} más")
    
    # Series
    series = chart_data.get('series', [])
    if series:
        description.append(f"Series de datos ({len(series)}):")
        for s in series:
            name = s.get('name', 'Sin nombre')
            values = s.get('values', [])
            if values:
                avg = sum(values) / len(values)
                description.append(f"   - {name}: {len(values)} valores, promedio: {avg:.1f}")
    
    return '\n'.join(description)


# Función de prueba
if __name__ == '__main__':
    import sys
    
    if len(sys.argv) < 2:
        print("Uso: python chart_modifier.py <archivo.pptx>")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    
    print(f"📊 Analizando gráficos en: {pptx_path}")
    
    prs = Presentation(pptx_path)
    chart_count = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_count += 1
                chart = shape.chart
                data = extract_chart_data(chart)
                
                print(f"\n--- Gráfico {chart_count} (Slide {slide_idx + 1}) ---")
                print(f"Tipo: {data.get('chart_type', 'Desconocido')}")
                print(f"Categorías: {data.get('categories', [])}")
                print(f"Series: {len(data.get('series', []))}")
                print(f"Título: {data.get('title', 'Sin título')}")
    
    print(f"\n📊 Total de gráficos encontrados: {chart_count}")
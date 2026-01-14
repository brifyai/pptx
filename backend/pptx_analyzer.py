from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Any
import json
import base64
from io import BytesIO
from PIL import Image
import os
import tempfile
import sys

# Agregar LibreOffice al path para UNO
LIBREOFFICE_PROGRAM = r"C:\Program Files\LibreOffice\program"
if os.path.exists(LIBREOFFICE_PROGRAM) and LIBREOFFICE_PROGRAM not in sys.path:
    sys.path.insert(0, LIBREOFFICE_PROGRAM)
    os.environ['PYTHONPATH'] = LIBREOFFICE_PROGRAM + os.pathsep + os.environ.get('PYTHONPATH', '')

# Importar procesador de imágenes
from image_processor import remove_white_background, smart_background_removal

# Intentar importar UNO API de LibreOffice
try:
    from libreoffice_uno_renderer import render_pptx_with_uno, UNO_AVAILABLE
    if UNO_AVAILABLE:
        print("✅ LibreOffice UNO API disponible - renderizado de alta calidad")
except ImportError as e:
    UNO_AVAILABLE = False
    print(f"⚠️ LibreOffice UNO API no disponible: {e}")

# Intentar importar conversores de imágenes
try:
    from pptx_full_renderer import render_pptx_complete
    FULL_RENDERER_AVAILABLE = True
except ImportError:
    FULL_RENDERER_AVAILABLE = False
    print("⚠️ Renderizador completo no disponible")

try:
    from pptx_renderer import render_pptx_to_images
    CUSTOM_RENDERER_AVAILABLE = True
except ImportError:
    CUSTOM_RENDERER_AVAILABLE = False
    print("⚠️ Renderizador personalizado no disponible")

try:
    from pptx_to_images_custom import convert_pptx_to_images_custom
    CUSTOM_RENDERER_OLD_AVAILABLE = True
except ImportError:
    CUSTOM_RENDERER_OLD_AVAILABLE = False

try:
    from pptx_to_images_aspose import convert_pptx_to_images_aspose
    ASPOSE_AVAILABLE = True
except ImportError:
    ASPOSE_AVAILABLE = False

try:
    from pptx_to_images import convert_pptx_to_images
    LIBREOFFICE_AVAILABLE = True
except ImportError:
    LIBREOFFICE_AVAILABLE = False

def analyze_presentation(pptx_path: str) -> Dict[str, Any]:
    """
    Analiza un archivo PowerPoint y extrae toda su estructura de diseño
    Incluye extracción de imágenes originales con transparencia
    """
    prs = Presentation(pptx_path)
    
    # Intentar generar previews de slides (prioridad: UNO API > LibreOffice > Full Renderer > Placeholder)
    slide_images = []
    
    if UNO_AVAILABLE:
        try:
            print("🎨 Usando LibreOffice UNO API (máxima calidad)...")
            slide_images = render_pptx_with_uno(pptx_path)
            print(f"✅ Generadas {len(slide_images)} imágenes con UNO API")
        except Exception as e:
            print(f"⚠️ Error con UNO API: {e}")
            import traceback
            traceback.print_exc()
    
    if not slide_images and LIBREOFFICE_AVAILABLE:
        try:
            print("🎨 Usando LibreOffice headless...")
            slide_images = convert_pptx_to_images(pptx_path)
            print(f"✅ Generadas {len(slide_images)} imágenes con LibreOffice")
        except Exception as e:
            print(f"⚠️ Error con LibreOffice: {e}")
            import traceback
            traceback.print_exc()
    
    if not slide_images and FULL_RENDERER_AVAILABLE:
        try:
            print("🎨 Usando renderizador completo...")
            slide_images = render_pptx_complete(pptx_path)
            print(f"✅ Generadas {len(slide_images)} imágenes con renderizador completo")
        except Exception as e:
            print(f"⚠️ Error con renderizador completo: {e}")
            import traceback
            traceback.print_exc()
    
    if not slide_images:
        print("⚠️ Usando placeholders")
        from pptx_to_images import generate_placeholder_image
        slide_images = [generate_placeholder_image(i + 1) for i in range(len(prs.slides))]
    
    # Extraer todos los assets (imágenes con transparencia, logos, etc.)
    print(f"\n{'='*60}")
    print(f"EXTRAYENDO ASSETS DEL PPTX")
    print(f"{'='*60}")
    extracted_assets = extract_all_assets(prs)
    print(f"\n📊 RESUMEN DE ASSETS:")
    print(f"   Total: {extracted_assets['totalCount']}")
    print(f"   Logos: {len(extracted_assets['logos'])}")
    print(f"   Transparentes: {len(extracted_assets['transparentImages'])}")
    print(f"   Animados: {len(extracted_assets['animatedElements'])}")
    print(f"   Imágenes: {len(extracted_assets['images'])}")
    print(f"{'='*60}\n")
    
    analysis = {
        "fileName": pptx_path.split('/')[-1].split('\\')[-1],
        "slideSize": {
            "width": prs.slide_width,
            "height": prs.slide_height
        },
        "slides": [],
        "slideImages": slide_images,
        "extractedAssets": extracted_assets,
        "renderMethod": "custom" if CUSTOM_RENDERER_AVAILABLE and slide_images else ("libreoffice" if LIBREOFFICE_AVAILABLE and slide_images else ("aspose" if ASPOSE_AVAILABLE and slide_images else "placeholder"))
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        # Extraer fondo del XML primero
        slide_bg = extract_background(slide)
        bg_color_hex = slide_bg.get('color', '#FFFFFF')
        
        # Si el fondo es blanco pero la preview muestra otro color, extraer de la imagen
        if bg_color_hex == '#FFFFFF' and slide_idx < len(slide_images):
            dominant_color = extract_dominant_color_from_preview(slide_images[slide_idx])
            # Solo usar el color dominante si NO es blanco
            if dominant_color != '#FFFFFF':
                print(f"   🎨 Color dominante detectado en preview: {dominant_color}")
                bg_color_hex = dominant_color
                slide_bg['color'] = dominant_color
        
        slide_data = {
            "number": slide_idx + 1,
            "type": detect_slide_type(slide),
            "layout": slide.slide_layout.name,
            "layoutType": get_layout_category(slide.slide_layout.name),  # Nuevo
            "isTitle": is_title_slide(slide),  # Nuevo
            "isCover": slide_idx == 0,  # Nuevo: primera slide es portada
            "background": slide_bg,
            "preview": slide_images[slide_idx] if slide_idx < len(slide_images) else None,
            "textAreas": [],
            "imageAreas": [],
            "shapes": []
        }
        
        # Analizar cada shape en la diapositiva
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_area = extract_text_area(shape)
                slide_data["textAreas"].append(text_area)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_area = extract_image_area(shape)
                slide_data["imageAreas"].append(image_area)
            
            else:
                shape_data = extract_shape_data(shape)
                slide_data["shapes"].append(shape_data)
        
        analysis["slides"].append(slide_data)
    
    return analysis

def detect_slide_type(slide) -> str:
    """
    Detecta el tipo de diapositiva basándose en múltiples criterios:
    - Layout name
    - Posición en la presentación
    - Cantidad y tipo de placeholders
    - Estructura de contenido
    """
    layout_name = slide.slide_layout.name.lower()
    slide_index = slide.slide_id  # Posición en la presentación
    
    # Contar placeholders por tipo
    title_count = 0
    subtitle_count = 0
    content_count = 0
    picture_count = 0
    
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # 1 = Title, 2 = Body/Content, 3 = Center Title, 13 = Subtitle, 18 = Picture
            if ph_type == 1 or ph_type == 3:  # Title or Center Title
                title_count += 1
            elif ph_type == 13:  # Subtitle
                subtitle_count += 1
            elif ph_type == 2:  # Body/Content
                content_count += 1
            elif ph_type == 18:  # Picture
                picture_count += 1
    
    # Detección por layout name (más específico)
    if 'title' in layout_name:
        if 'only' in layout_name or 'slide' in layout_name:
            return 'title'  # Portada
        elif 'section' in layout_name:
            return 'section'  # Separador de sección
        elif 'content' in layout_name:
            return 'title_content'  # Título con contenido
    
    if 'blank' in layout_name or 'en blanco' in layout_name:
        return 'blank'
    
    if 'section' in layout_name or 'sección' in layout_name:
        return 'section'
    
    if 'comparison' in layout_name or 'comparación' in layout_name:
        return 'comparison'
    
    if 'two' in layout_name or 'dos' in layout_name:
        return 'two_content'
    
    if 'picture' in layout_name or 'imagen' in layout_name:
        return 'picture'
    
    if 'quote' in layout_name or 'cita' in layout_name:
        return 'quote'
    
    # Detección por estructura de placeholders
    if title_count > 0 and subtitle_count > 0 and content_count == 0:
        return 'title'  # Portada típica: título + subtítulo
    
    if title_count > 0 and content_count > 0:
        return 'content'  # Slide de contenido típico
    
    if title_count == 0 and content_count == 0 and picture_count > 0:
        return 'picture'  # Solo imágenes
    
    if title_count > 0 and content_count == 0 and subtitle_count == 0:
        return 'section'  # Separador de sección
    
    # Detección por posición (heurística)
    # La primera slide suele ser portada
    if slide_index == 1 and title_count > 0:
        return 'title'
    
    # Default
    return 'content'


def get_layout_category(layout_name: str) -> str:
    """
    Categoriza el layout en tipos generales
    """
    layout_lower = layout_name.lower()
    
    if 'title' in layout_lower and ('only' in layout_lower or 'slide' in layout_lower):
        return 'cover'
    elif 'title' in layout_lower:
        return 'title_slide'
    elif 'section' in layout_lower or 'sección' in layout_lower:
        return 'section_header'
    elif 'blank' in layout_lower or 'en blanco' in layout_lower:
        return 'blank'
    elif 'two' in layout_lower or 'dos' in layout_lower or 'comparison' in layout_lower:
        return 'two_column'
    elif 'picture' in layout_lower or 'imagen' in layout_lower:
        return 'picture_focused'
    elif 'content' in layout_lower or 'contenido' in layout_lower:
        return 'content'
    else:
        return 'other'


def is_title_slide(slide) -> bool:
    """
    Determina si una slide es de tipo título/portada
    basándose en la estructura de placeholders
    """
    has_title = False
    has_subtitle = False
    has_content = False
    
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in [1, 3]:  # Title or Center Title
                has_title = True
            elif ph_type == 13:  # Subtitle
                has_subtitle = True
            elif ph_type == 2:  # Body/Content
                has_content = True
    
    # Es título si tiene título y subtítulo pero NO contenido
    return has_title and has_subtitle and not has_content

def get_theme_colors(slide) -> Dict[str, str]:
    """
    Extrae los colores del tema del PPTX
    """
    try:
        from lxml import etree
        
        namespaces = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        
        # Acceder al tema a través del master
        master = slide.slide_layout.slide_master
        theme_part = master.part.part_related_by('http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
        theme_xml = theme_part.blob
        theme_root = etree.fromstring(theme_xml)
        
        # Buscar el esquema de colores
        color_scheme = theme_root.find('.//a:clrScheme', namespaces)
        
        if color_scheme is None:
            return {}
        
        theme_colors = {}
        
        for color_elem in color_scheme:
            color_name = color_elem.tag.split('}')[-1]
            
            # Buscar el valor del color
            srgb = color_elem.find('.//a:srgbClr', namespaces)
            sys_clr = color_elem.find('.//a:sysClr', namespaces)
            
            if srgb is not None:
                color_val = srgb.get('val')
                theme_colors[color_name] = f"#{color_val}"
            elif sys_clr is not None:
                color_val = sys_clr.get('lastClr')
                if color_val:
                    theme_colors[color_name] = f"#{color_val}"
        
        return theme_colors
        
    except Exception as e:
        print(f"   ⚠️ Error extrayendo colores del tema: {e}")
        return {}


def extract_background(slide) -> Dict[str, Any]:
    """
    Extrae información del fondo de la diapositiva
    Lee directamente del XML y usa los colores reales del tema
    
    Estrategia mejorada:
    1. Buscar fondo en el slide XML
    2. Si no hay, buscar en el layout
    3. Si no hay, buscar en el slide master
    4. Si no hay, usar color del tema (lt1 o bg1)
    5. Si todo falla, usar el color dominante del preview
    """
    background = {
        "type": "solid",
        "color": None,
        "source": "unknown"  # slide, layout, master, theme, preview
    }
    
    try:
        from lxml import etree
        
        # Obtener colores del tema primero
        theme_colors = get_theme_colors(slide)
        
        # Mapeo de colores del tema a nombres legibles
        theme_color_names = {
            'bg1': 'Fondo 1',
            'bg2': 'Fondo 2',
            'tx1': 'Texto 1',
            'tx2': 'Texto 2',
            'lt1': 'Claro 1',
            'lt2': 'Claro 2',
            'dk1': 'Oscuro 1',
            'dk2': 'Oscuro 2',
            'accent1': 'Acento 1',
            'accent2': 'Acento 2',
            'accent3': 'Acento 3',
            'accent4': 'Acento 4',
            'accent5': 'Acento 5',
            'accent6': 'Acento 6',
        }
        
        def apply_theme_color(scheme_val, source_name):
            """Aplica un color del tema al fondo"""
            if scheme_val in theme_colors:
                background["color"] = theme_colors[scheme_val]
                background["source"] = source_name
                color_name = theme_color_names.get(scheme_val, scheme_val)
                print(f"   ✅ Color del tema ({color_name}) aplicado desde {source_name}: {background['color']}")
            else:
                # Fallback a colores por defecto
                scheme_colors_default = {
                    'bg1': '#FFFFFF',
                    'bg2': '#F2F2F2',
                    'tx1': '#000000',
                    'tx2': '#1F1F1F',
                    'accent1': '#4472C4',
                    'accent2': '#ED7D31',
                    'accent3': '#A5A5A5',
                    'accent4': '#FFC000',
                    'accent5': '#5B9BD5',
                    'accent6': '#70AD47',
                    'dk1': '#000000',
                    'lt1': '#FFFFFF',
                    'dk2': '#1F1F1F',
                    'lt2': '#EEECE1'
                }
                background["color"] = scheme_colors_default.get(scheme_val, '#FFFFFF')
                background["source"] = f"{source_name}_default"
                color_name = theme_color_names.get(scheme_val, scheme_val)
                print(f"   ⚠️ Color por defecto ({color_name}) aplicado: {background['color']}")
        
        # Namespaces de PowerPoint
        namespaces = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        
        # ===== 1. Buscar en el slide XML =====
        slide_part = slide.part
        slide_xml = slide_part.blob
        root = etree.fromstring(slide_xml)
        
        bg_element = root.find('.//p:cSld/p:bg', namespaces)
        
        if bg_element is not None:
            # Color sólido directo
            solid_fill = bg_element.find('.//a:solidFill/a:srgbClr', namespaces)
            if solid_fill is not None:
                color_val = solid_fill.get('val')
                if color_val:
                    background["color"] = f"#{color_val}"
                    background["type"] = "solid"
                    background["source"] = "slide"
                    print(f"   ✅ Color de fondo extraído del slide: {background['color']}")
                    return background
            
            # Color del esquema
            scheme_fill = bg_element.find('.//a:solidFill/a:schemeClr', namespaces)
            if scheme_fill is not None:
                scheme_val = scheme_fill.get('val')
                apply_theme_color(scheme_val, "slide")
                return background
            
            # Gradiente
            grad_fill = bg_element.find('.//a:gradFill', namespaces)
            if grad_fill is not None:
                background["type"] = "gradient"
                first_color = grad_fill.find('.//a:srgbClr', namespaces)
                if first_color is not None:
                    color_val = first_color.get('val')
                    if color_val:
                        background["color"] = f"#{color_val}"
                        background["source"] = "slide_gradient"
                        print(f"   ✅ Color de gradiente extraído del slide: {background['color']}")
                        return background
        
        # ===== 2. Buscar en el layout =====
        print(f"   ℹ️ No se encontró fondo en el slide, buscando en layout...")
        layout_part = slide.slide_layout.part
        layout_xml = layout_part.blob
        layout_root = etree.fromstring(layout_xml)
        
        bg_element = layout_root.find('.//p:cSld/p:bg', namespaces)
        if bg_element is not None:
            solid_fill = bg_element.find('.//a:solidFill/a:srgbClr', namespaces)
            if solid_fill is not None:
                color_val = solid_fill.get('val')
                if color_val:
                    background["color"] = f"#{color_val}"
                    background["type"] = "solid"
                    background["source"] = "layout"
                    print(f"   ✅ Color de fondo extraído del layout: {background['color']}")
                    return background
            
            scheme_fill = bg_element.find('.//a:solidFill/a:schemeClr', namespaces)
            if scheme_fill is not None:
                scheme_val = scheme_fill.get('val')
                apply_theme_color(scheme_val, "layout")
                return background
        
        # ===== 3. Buscar en el slide master =====
        print(f"   ℹ️ No se encontró fondo en el layout, buscando en master...")
        try:
            master = slide.slide_layout.slide_master
            master_part = master.part
            master_xml = master_part.blob
            master_root = etree.fromstring(master_xml)
            
            bg_element = master_root.find('.//p:cSld/p:bg', namespaces)
            if bg_element is not None:
                solid_fill = bg_element.find('.//a:solidFill/a:srgbClr', namespaces)
                if solid_fill is not None:
                    color_val = solid_fill.get('val')
                    if color_val:
                        background["color"] = f"#{color_val}"
                        background["type"] = "solid"
                        background["source"] = "master"
                        print(f"   ✅ Color de fondo extraído del master: {background['color']}")
                        return background
                
                scheme_fill = bg_element.find('.//a:solidFill/a:schemeClr', namespaces)
                if scheme_fill is not None:
                    scheme_val = scheme_fill.get('val')
                    apply_theme_color(scheme_val, "master")
                    return background
        except Exception as master_error:
            print(f"   ⚠️ Error buscando en master: {master_error}")
        
        # ===== 4. Usar color del tema directamente =====
        print(f"   ℹ️ No se encontró fondo explícito, usando color del tema...")
        
        # Preferir lt1 (usual para fondos claros) o bg1
        if 'lt1' in theme_colors:
            background["color"] = theme_colors['lt1']
            background["source"] = "theme_lt1"
            print(f"   ✅ Color del tema (lt1) aplicado: {background['color']}")
        elif 'bg1' in theme_colors:
            background["color"] = theme_colors['bg1']
            background["source"] = "theme_bg1"
            print(f"   ✅ Color del tema (bg1) aplicado: {background['color']}")
        elif theme_colors:
            # Usar el primer color del tema disponible
            first_color = list(theme_colors.values())[0]
            background["color"] = first_color
            background["source"] = "theme_first"
            print(f"   ✅ Primer color del tema aplicado: {background['color']}")
        
    except Exception as e:
        print(f"   ⚠️ Error extrayendo fondo del XML: {e}")
        import traceback
        traceback.print_exc()
    
    # Si todo falla, #FFFFFF por defecto
    if not background["color"]:
        background["color"] = "#FFFFFF"
        background["source"] = "default_white"
        print(f"   ℹ️ No se pudo detectar color de fondo, usando blanco por defecto")
    
    return background

def extract_text_area(shape) -> Dict[str, Any]:
    """
    Extrae información completa de un área de texto con detalles de diseño
    """
    text_frame = shape.text_frame
    
    # Detectar tipo de texto
    text_type = "body"
    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type == 1:  # Title
            text_type = "title"
        elif placeholder_type == 2:  # Subtitle
            text_type = "subtitle"
        elif placeholder_type == 7:  # Body
            text_type = "bullets"
    
    # Extraer formato detallado
    formatting = extract_text_formatting(text_frame)
    
    text_area = {
        "id": shape.shape_id,
        "type": text_type,
        "position": {
            "x": shape.left,
            "y": shape.top,
            "width": shape.width,
            "height": shape.height,
            "x_percent": (shape.left / 9144000) * 100,  # Convertir a porcentaje
            "y_percent": (shape.top / 6858000) * 100,
            "width_percent": (shape.width / 9144000) * 100,
            "height_percent": (shape.height / 6858000) * 100
        },
        "text": text_frame.text,
        "formatting": formatting,
        "maxChars": estimate_max_chars(shape.width, shape.height, formatting.get('size', 18)),
        "alignment": str(text_frame.paragraphs[0].alignment) if len(text_frame.paragraphs) > 0 else None,
        "canEdit": True
    }
    
    return text_area


def estimate_max_chars(width, height, font_size):
    """
    Estima el número máximo de caracteres que caben en un área de texto
    """
    # Asegurar que font_size no sea None
    if font_size is None:
        font_size = 18
    
    # Asegurar que width y height no sean None
    if width is None or height is None:
        return 500  # Valor por defecto
    
    # Aproximación: ancho en EMUs / (font_size * 9525) para caracteres por línea
    # altura en EMUs / (font_size * 12700) para número de líneas
    chars_per_line = max(1, int(width / (font_size * 9525)))
    num_lines = max(1, int(height / (font_size * 12700)))
    return chars_per_line * num_lines

def extract_text_formatting(text_frame) -> Dict[str, Any]:
    """
    Extrae el formato del texto (fuente, tamaño, color, etc.)
    """
    formatting = {
        "font": None,
        "size": None,
        "color": None,
        "bold": False,
        "italic": False,
        "alignment": None
    }
    
    if len(text_frame.paragraphs) > 0:
        paragraph = text_frame.paragraphs[0]
        if len(paragraph.runs) > 0:
            run = paragraph.runs[0]
            font = run.font
            
            formatting["font"] = font.name
            formatting["size"] = font.size.pt if font.size else None
            formatting["bold"] = font.bold
            formatting["italic"] = font.italic
            
            # Acceso seguro a color.rgb
            try:
                if font.color and hasattr(font.color, 'rgb') and font.color.rgb is not None:
                    formatting["color"] = rgb_to_hex(font.color.rgb)
            except (AttributeError, TypeError):
                pass
            
            formatting["alignment"] = str(paragraph.alignment)
    
    return formatting

def extract_image_area(shape) -> Dict[str, Any]:
    """
    Extrae información de un área de imagen, incluyendo la imagen original con transparencia
    """
    image_data = {
        "id": shape.shape_id,
        "position": {
            "x": shape.left,
            "y": shape.top,
            "width": shape.width,
            "height": shape.height
        },
        "type": "image",
        "imageBase64": None,
        "hasTransparency": False,
        "format": None,
        "isLogo": False
    }
    
    try:
        # Extraer la imagen original del shape
        image = shape.image
        image_bytes = image.blob
        content_type = image.content_type
        
        # Determinar formato
        if 'png' in content_type:
            image_data["format"] = "png"
        elif 'jpeg' in content_type or 'jpg' in content_type:
            image_data["format"] = "jpeg"
        elif 'gif' in content_type:
            image_data["format"] = "gif"
        elif 'svg' in content_type:
            image_data["format"] = "svg"
        else:
            image_data["format"] = content_type.split('/')[-1] if '/' in content_type else "unknown"
        
        # Verificar si tiene transparencia (PNG con canal alpha)
        if image_data["format"] == "png":
            try:
                img = Image.open(BytesIO(image_bytes))
                if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                    image_data["hasTransparency"] = True
                    print(f"🔍 Imagen con transparencia detectada: shape_id={shape.shape_id}")
            except Exception as e:
                print(f"⚠️ Error verificando transparencia: {e}")
        
        # Detectar si es un logo (imagen pequeña, posiblemente en esquina)
        slide_width = 9144000  # Ancho típico de slide en EMUs
        slide_height = 6858000  # Alto típico
        
        img_width_ratio = shape.width / slide_width
        img_height_ratio = shape.height / slide_height
        
        # Si la imagen ocupa menos del 25% del slide, podría ser un logo
        if img_width_ratio < 0.25 and img_height_ratio < 0.25:
            image_data["isLogo"] = True
            print(f"🏷️ Posible logo detectado: shape_id={shape.shape_id}")
        
        # Convertir a base64 para enviar al frontend
        img_base64 = base64.b64encode(image_bytes).decode('utf-8')
        mime_type = content_type if content_type else f"image/{image_data['format']}"
        image_data["imageBase64"] = f"data:{mime_type};base64,{img_base64}"
        
        print(f"✅ Imagen extraída: format={image_data['format']}, transparency={image_data['hasTransparency']}, logo={image_data['isLogo']}")
        
    except Exception as e:
        print(f"⚠️ No se pudo extraer imagen del shape {shape.shape_id}: {e}")
    
    return image_data

def extract_shape_data(shape) -> Dict[str, Any]:
    """
    Extrae información de formas (rectángulos, círculos, etc.)
    """
    return {
        "id": shape.shape_id,
        "type": str(shape.shape_type),
        "position": {
            "x": shape.left,
            "y": shape.top,
            "width": shape.width,
            "height": shape.height
        }
    }

def get_color_rgb(color):
    """
    Obtiene el color RGB de un objeto de color de forma segura
    """
    try:
        if color and hasattr(color, 'rgb') and color.rgb is not None:
            return rgb_to_hex(color.rgb)
    except (AttributeError, TypeError):
        pass
    return None

def rgb_to_hex(rgb):
    """
    Convierte RGB a hexadecimal
    """
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


def extract_dominant_color_from_preview(preview_base64: str) -> str:
    """
    Extrae el color dominante de una imagen preview del slide
    Útil cuando el fondo no se puede extraer del XML
    """
    try:
        from PIL import Image
        import base64
        from io import BytesIO
        from collections import Counter
        
        # Decodificar imagen
        if ',' in preview_base64:
            preview_base64 = preview_base64.split(',')[1]
        
        image_data = base64.b64decode(preview_base64)
        img = Image.open(BytesIO(image_data))
        
        # Redimensionar para análisis más rápido
        img = img.resize((100, 100))
        img = img.convert('RGB')
        
        # Obtener pixels de las esquinas (donde suele estar el fondo)
        pixels = []
        width, height = img.size
        
        # Muestrear esquinas y bordes
        for x in range(0, width, 10):
            pixels.append(img.getpixel((x, 0)))  # Borde superior
            pixels.append(img.getpixel((x, height-1)))  # Borde inferior
        
        for y in range(0, height, 10):
            pixels.append(img.getpixel((0, y)))  # Borde izquierdo
            pixels.append(img.getpixel((width-1, y)))  # Borde derecho
        
        # Encontrar el color más común
        color_counts = Counter(pixels)
        dominant_color = color_counts.most_common(1)[0][0]
        
        # Convertir a hex
        hex_color = f"#{dominant_color[0]:02x}{dominant_color[1]:02x}{dominant_color[2]:02x}"
        
        return hex_color
        
    except Exception as e:
        print(f"   ⚠️ Error extrayendo color dominante: {e}")
        return "#FFFFFF"


def extract_all_assets(prs) -> Dict[str, Any]:
    """
    Extrae todos los assets (imágenes, logos) del PPTX
    Preserva transparencias y formatos originales
    Incluye el color de fondo del slide para cada asset
    """
    assets = {
        "logos": [],
        "images": [],
        "transparentImages": [],
        "animatedElements": [],  # Elementos con animación
        "totalCount": 0
    }
    
    seen_blobs = set()  # Para evitar duplicados
    
    for slide_idx, slide in enumerate(prs.slides):
        # Detectar elementos animados en este slide
        print(f"\n🔍 Analizando slide {slide_idx + 1} para detectar animaciones...")
        animated_shape_ids = detect_animated_shapes(slide, prs)
        print(f"   Resultado: {len(animated_shape_ids)} shapes animados detectados: {animated_shape_ids}")
        
        # Extraer color de fondo del slide
        slide_bg_color = extract_background(slide)
        bg_color_hex = slide_bg_color.get('color', '#FFFFFF')  # Default blanco
        print(f"   🎨 Color de fondo del slide: {bg_color_hex}")
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    blob_hash = hash(image.blob[:100])  # Hash parcial para detectar duplicados
                    
                    if blob_hash in seen_blobs:
                        continue
                    seen_blobs.add(blob_hash)
                    
                    image_bytes = image.blob
                    content_type = image.content_type
                    
                    # Determinar formato
                    img_format = "png"
                    if 'jpeg' in content_type or 'jpg' in content_type:
                        img_format = "jpeg"
                    elif 'gif' in content_type:
                        img_format = "gif"
                    elif 'png' in content_type:
                        img_format = "png"
                    
                    # Verificar transparencia
                    has_transparency = False
                    if img_format == "png":
                        try:
                            img = Image.open(BytesIO(image_bytes))
                            if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                                has_transparency = True
                        except:
                            pass
                    
                    # Detectar si es logo
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    is_logo = (shape.width / slide_width < 0.25) and (shape.height / slide_height < 0.25)
                    
                    # Detectar si tiene animación
                    has_animation = shape.shape_id in animated_shape_ids
                    
                    # Crear asset
                    img_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    mime_type = content_type if content_type else f"image/{img_format}"
                    original_image = f"data:{mime_type};base64,{img_base64}"
                    
                    # PROCESAR IMAGEN: Remover fondo blanco y aplicar color del slide
                    processed_image = original_image
                    if bg_color_hex and bg_color_hex != "#FFFFFF":
                        print(f"      🎨 Procesando imagen para aplicar fondo {bg_color_hex}...")
                        processed_image = smart_background_removal(original_image, bg_color_hex)
                        print(f"      ✅ Imagen procesada con nuevo fondo")
                    
                    asset = {
                        "id": f"asset_{slide_idx}_{shape.shape_id}",
                        "slideNumber": slide_idx + 1,
                        "shapeId": shape.shape_id,
                        "format": img_format,
                        "hasTransparency": has_transparency,
                        "hasAnimation": has_animation,
                        "isLogo": is_logo,
                        "backgroundColor": bg_color_hex,  # Color de fondo del slide
                        "position": {
                            "x": shape.left,
                            "y": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        },
                        "imageBase64": processed_image  # Usar imagen procesada
                    }
                    
                    # Clasificar - priorizar elementos animados
                    if has_animation:
                        assets["animatedElements"].append(asset)
                        print(f"🎬 Elemento animado extraído: slide {slide_idx + 1}, shape_id={shape.shape_id}")
                    elif is_logo:
                        assets["logos"].append(asset)
                        print(f"🏷️ Logo extraído: slide {slide_idx + 1}, transparency={has_transparency}")
                    elif has_transparency:
                        assets["transparentImages"].append(asset)
                        print(f"🔍 Imagen transparente extraída: slide {slide_idx + 1}")
                    else:
                        assets["images"].append(asset)
                    
                    assets["totalCount"] += 1
                    
                except Exception as e:
                    print(f"⚠️ Error extrayendo asset: {e}")
    
    print(f"📦 Assets extraídos: {assets['totalCount']} total ({len(assets['logos'])} logos, {len(assets['transparentImages'])} transparentes, {len(assets['animatedElements'])} animados, {len(assets['images'])} imágenes)")
    
    return assets


def detect_animated_shapes(slide, prs) -> set:
    """
    Detecta qué shapes tienen animación en un slide.
    Lee directamente el XML del slide para encontrar elementos p:timing.
    ESTRATEGIA: Si no se encuentran animaciones en el XML, asumir que las imágenes
    pequeñas con transparencia (logos) probablemente tienen animación.
    """
    animated_ids = set()
    
    try:
        # Acceder al XML del slide a través de la parte del slide
        from lxml import etree
        
        # Obtener el XML del slide
        slide_part = slide.part
        
        # Leer el XML directamente
        slide_xml = slide_part.blob
        root = etree.fromstring(slide_xml)
        
        # Namespaces de PowerPoint - probar múltiples variantes
        namespaces_variants = [
            {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            },
            {
                'p': 'http://schemas.microsoft.com/office/powerpoint/2006/main',
                'a': 'http://schemas.microsoft.com/office/drawing/2006/main',
            }
        ]
        
        timing_found = False
        
        # Intentar con diferentes namespaces
        for namespaces in namespaces_variants:
            # Buscar elementos de timing/animación
            timing = root.find('.//p:timing', namespaces)
            
            if timing is not None:
                timing_found = True
                print(f"   🎬 Slide tiene elemento p:timing (namespace: {namespaces['p']})")
                
                # Buscar todos los targetElement que referencian shapes
                # Pueden estar en diferentes estructuras
                targets = timing.findall('.//*[@spid]', namespaces)
                
                for target in targets:
                    shape_id = target.get('spid')
                    if shape_id:
                        try:
                            animated_ids.add(int(shape_id))
                            print(f"   🎬 Shape {shape_id} tiene animación")
                        except ValueError:
                            pass
                
                # También buscar en p:tgtEl/p:spTgt
                sp_targets = timing.findall('.//p:tgtEl/p:spTgt', namespaces)
                for target in sp_targets:
                    shape_id = target.get('spid')
                    if shape_id:
                        try:
                            animated_ids.add(int(shape_id))
                            print(f"   🎬 Shape {shape_id} tiene animación (spTgt)")
                        except ValueError:
                            pass
                
                break  # Si encontramos timing, no necesitamos probar otros namespaces
        
        # Si no encontramos timing con namespaces, buscar iterando por todos los elementos
        if not timing_found:
            # Buscar manualmente sin namespace
            for elem in root.iter():
                tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                if tag_name == 'timing':
                    timing_found = True
                    print(f"   🎬 Slide tiene elemento timing (encontrado sin namespace)")
                    
                    # Buscar targets en este elemento
                    for target in elem.iter():
                        shape_id = target.get('spid')
                        if shape_id:
                            try:
                                animated_ids.add(int(shape_id))
                                print(f"   🎬 Shape {shape_id} tiene animación")
                            except ValueError:
                                pass
                    break
        
        if not timing_found:
            print(f"   ℹ️ Slide no tiene animaciones detectadas en XML")
            print(f"   🔍 Aplicando FALLBACK MEJORADO: detectar solo elementos con alta probabilidad de animación...")
            
            # FALLBACK MEJORADO: Solo marcar elementos con alta probabilidad de animación
            # NO marcar todas las imágenes - solo las que cumplen criterios estrictos
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from PIL import Image
            from io import BytesIO
            
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            print(f"   📏 Dimensiones del slide: {slide_width} x {slide_height} EMUs")
            
            picture_count = 0
            high_prob_animations = 0
            medium_prob_animations = 0
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_count += 1
                    try:
                        # Verificar criterios para animación
                        width_ratio = shape.width / slide_width
                        height_ratio = shape.height / slide_height
                        is_small = width_ratio < 0.15 and height_ratio < 0.15  # Más estricto: < 15%
                        
                        print(f"   📷 Imagen encontrada: shape_id={shape.shape_id}")
                        print(f"      Tamaño: {width_ratio:.1%} x {height_ratio:.1%} del slide")
                        print(f"      ¿Es pequeña (<15%)? {is_small}")
                        
                        # Verificar si tiene transparencia
                        image = shape.image
                        image_bytes = image.blob
                        content_type = image.content_type
                        
                        print(f"      Formato: {content_type}")
                        
                        has_transparency = False
                        if 'png' in content_type:
                            try:
                                img = Image.open(BytesIO(image_bytes))
                                print(f"      Modo de imagen: {img.mode}")
                                if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                                    has_transparency = True
                                    print(f"      ✅ Tiene transparencia!")
                            except Exception as e:
                                print(f"      ⚠️ Error verificando transparencia: {e}")
                        
                        print(f"      ¿Tiene transparencia? {has_transparency}")
                        
                        # CRITERIOS MEJORADOS para marcar como animado:
                        # ALTA PROBABILIDAD: imagen pequeña (<15%) + transparencia + posición en esquina
                        # MEDIA PROBABILIDAD: imagen pequeña (<15%) + transparencia (sin importar posición)
                        # NO ANIMADA: imágenes grandes o sin transparencia
                        
                        # Verificar si está en una esquina (común para logos animados)
                        is_corner = (
                            (shape.left < slide_width * 0.1 or shape.left > slide_width * 0.85) and
                            (shape.top < slide_height * 0.1 or shape.top > slide_height * 0.85)
                        )
                        print(f"      ¿Está en esquina? {is_corner}")
                        
                        if is_small and has_transparency:
                            if is_corner:
                                # ALTA PROBABILIDAD - logo en esquina con transparencia
                                animated_ids.add(shape.shape_id)
                                high_prob_animations += 1
                                print(f"   🎬 ✅ Shape {shape.shape_id} marcado como ALTA probabilidad (logo en esquina con transparencia)")
                            else:
                                # MEDIA PROBABILIDAD - imagen pequeña con transparencia
                                animated_ids.add(shape.shape_id)
                                medium_prob_animations += 1
                                print(f"   🎬 ⚠️ Shape {shape.shape_id} marcado como MEDIA PROBABILIDAD (imagen pequeña con transparencia)")
                        else:
                            # BAJA PROBABILIDAD - no marcar como animado
                            print(f"   🎬 ❌ Shape {shape.shape_id} NO marcado (baja probabilidad de animación)")
                        
                    except Exception as e:
                        print(f"      ⚠️ Error procesando imagen: {e}")
            
            print(f"   📊 Total de imágenes analizadas: {picture_count}")
            print(f"   🎬 Animaciones detectadas: ALTA={high_prob_animations}, MEDIA={medium_prob_animations}")
            print(f"   ℹ️ Total marcado como animado: {len(animated_ids)}")
            print(f"   ℹ️ NOTA: Solo se marcan elementos con ALTA/MEDIA probabilidad de animación")
        
    except Exception as e:
        print(f"⚠️ Error detectando animaciones: {e}")
        import traceback
        traceback.print_exc()
    
    return animated_ids

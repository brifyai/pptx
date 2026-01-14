from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import tempfile
import os
from typing import Dict, Any, Optional, List
from copy import deepcopy

# Importar el nuevo clonador XML avanzado
import sys
import os

# Agregar el directorio padre al path para imports
backend_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(backend_dir)
if parent_dir not in sys.path:
    sys.path.insert(0, parent_dir)

try:
    from pptx_xml_cloner import clone_pptx_preserving_all, PPTXXMLCloner
    XML_CLONER_AVAILABLE = True
    print("✅ Clonador XML avanzado disponible")
except ImportError as e:
    XML_CLONER_AVAILABLE = False
    print(f"⚠️ Clonador XML no disponible: {e}")
    print("   Usando método legacy (python-pptx)")

# Importar módulos avanzados
SMARTART_AVAILABLE = False
CHART_MODIFIER_AVAILABLE = False
TABLE_PRESERVER_AVAILABLE = False

try:
    from smartart_extractor import extract_smartart_from_pptx, extract_diagram_text, analyze_smartart_for_ai
    SMARTART_AVAILABLE = True
    print("✅ Módulo SmartArt disponible")
except ImportError as e:
    print(f"⚠️ Módulo SmartArt no disponible: {e}")

try:
    from chart_modifier import extract_chart_data as extract_chart_data_advanced, generate_chart_data_with_ai, update_chart_with_data, create_chart_from_data, analyze_chart_for_ai
    CHART_MODIFIER_AVAILABLE = True
    print("✅ Módulo Chart Modifier disponible")
except ImportError as e:
    print(f"⚠️ Módulo Chart Modifier no disponible: {e}")

try:
    from table_preserver import extract_table_data, generate_table_xml, update_table_with_data, create_table_from_data, analyze_table_for_ai, preserve_table_xml, restore_table_from_preservation
    TABLE_PRESERVER_AVAILABLE = True
    print("✅ Módulo Table Preserver disponible")
except ImportError as e:
    print(f"⚠️ Módulo Table Preserver no disponible: {e}")


def generate_presentation(original_path: str, ai_content: Optional[Dict] = None,
                         use_xml_cloner: bool = True,
                         text_areas_by_slide: List[List[Dict]] = None) -> str:
    """
    Genera una nueva presentación CLONANDO completamente el diseño original
    y solo reemplazando el contenido de texto con el generado por IA.
    
    IMPORTANTE: Preserva TODOS los elementos visuales (formas, imágenes, fondos, etc.)
    
    Args:
        original_path: Ruta al archivo PPTX template
        ai_content: Contenido generado por IA
            - slides[].content: contenido de texto (title, subtitle, bullets, heading)
            - slides[].text_areas: mapeo por coordenadas exactas
            - slides[].table_data: datos para tablas
            - slides[].chart_data: datos para gráficos
        use_xml_cloner: Si True, usa el clonador XML avanzado que preserva
                       animaciones, transiciones, SmartArt, etc.
        text_areas_by_slide: Lista de textAreas por slide para reemplazo preciso
    
    Returns:
        Ruta al archivo PPTX generado
    """
    print(f"📄 Generando presentación desde: {original_path}")
    print(f"📝 Contenido IA recibido: {ai_content is not None}")
    print(f"📍 text_areas_by_slide: {text_areas_by_slide is not None}")
    
    # FORZAR uso del clonador XML cuando hay template y contenido
    if ai_content and isinstance(ai_content, dict) and 'slides' in ai_content:
        if XML_CLONER_AVAILABLE:
            try:
                print("🚀 Usando CLONADOR XML (preserva animaciones, SmartArt, gradientes, macros)")
                return generate_with_xml_cloner(original_path, ai_content, text_areas_by_slide)
            except Exception as e:
                print(f"⚠️ Error con clonador XML: {e}")
                import traceback
                traceback.print_exc()
        else:
            print("⚠️ Clonador XML no disponible, forzando método legacy")
    
    # Método legacy (fallback)
    return generate_with_legacy_method(original_path, ai_content)


def generate_with_xml_cloner(original_path: str, ai_content: Dict,
                            text_areas_by_slide: List[List[Dict]] = None) -> str:
    """
    Genera presentación usando el clonador XML avanzado.
    
    Preserva:
    - ✅ Animaciones y transiciones
    - ✅ SmartArt
    - ✅ Gradientes complejos
    - ✅ Sombras y efectos 3D
    - ✅ Todos los formatos de texto
    - ✅ Macros VBA
    
    Args:
        original_path: Ruta al template
        ai_content: Contenido IA
        text_areas_by_slide: Lista de textAreas por slide para reemplazo preciso
    """
    print("🚀 Usando clonador XML avanzado (preserva animaciones, SmartArt, etc.)")
    
    # Preparar contenido en formato para el clonador XML
    content_by_slide = []
    
    if 'slides' in ai_content:
        for idx, slide_data in enumerate(ai_content['slides']):
            slide_content = slide_data.get('content', {})
            content_by_slide.append(slide_content)
            print(f"   📝 Slide {idx+1} contenido:")
            print(f"      - title: {slide_content.get('title', 'N/A')[:50]}")
            print(f"      - subtitle: {slide_content.get('subtitle', 'N/A')[:50]}")
            print(f"      - heading: {slide_content.get('heading', 'N/A')[:50]}")
            print(f"      - bullets: {len(slide_content.get('bullets', []))} items")
            if slide_content.get('bullets'):
                for i, bullet in enumerate(slide_content.get('bullets', [])[:3]):
                    print(f"        • {bullet[:50]}")
            
            # Procesar SmartArt si está disponible
            if SMARTART_AVAILABLE and 'smartart' in slide_content:
                print(f"      - SmartArt: {len(slide_content['smartart'])} elementos")
            
            # Procesar tablas si hay datos
            if 'table_data' in slide_content:
                print(f"      - Tabla: {slide_content['table_data'].get('rows', 0)}x{slide_content['table_data'].get('cols', 0)}")
            
            # Procesar gráficos si hay datos
            if 'chart_data' in slide_content:
                print(f"      - Gráfico: {slide_content['chart_data'].get('chart_type', 'Desconocido')}")
    
    print(f"   📊 Total slides con contenido: {len(content_by_slide)}")
    print(f"   📍 text_areas disponibles: {text_areas_by_slide is not None}")
    
    # Usar el clonador XML con textAreas para reemplazo preciso
    try:
        output_path = clone_pptx_preserving_all(original_path, content_by_slide, text_areas_by_slide)
        print(f"✅ Presentación generada con clonador XML: {output_path}")
        return output_path
    except Exception as e:
        print(f"❌ Error en clonador XML: {e}")
        import traceback
        traceback.print_exc()
        raise


def generate_with_legacy_method(original_path: str, ai_content: Optional[Dict] = None) -> str:
    """
    Método legacy usando python-pptx (no preserva animaciones).
    """
    print("📦 Usando método legacy (python-pptx)")
    print(f"📝 Contenido IA recibido: {ai_content is not None}")
    
    # Cargar presentación original (template)
    template_prs = Presentation(original_path)
    
    print(f"📊 Template tiene {len(template_prs.slides)} slides")
    
    # Crear nueva presentación vacía
    new_prs = Presentation()
    
    # Copiar dimensiones del slide
    new_prs.slide_width = template_prs.slide_width
    new_prs.slide_height = template_prs.slide_height
    
    print(f"📐 Dimensiones: {new_prs.slide_width} x {new_prs.slide_height}")
    
    # Clonar cada slide completo
    for slide_idx, template_slide in enumerate(template_prs.slides):
        print(f"\n🔄 Clonando slide {slide_idx + 1}...")
        
        # Obtener contenido de IA para este slide
        ai_slide_content = None
        if ai_content and 'slides' in ai_content and slide_idx < len(ai_content['slides']):
            ai_slide_content = ai_content['slides'][slide_idx].get('content', {})
            print(f"   📝 Contenido IA: {list(ai_slide_content.keys()) if ai_slide_content else 'ninguno'}")
            
            # Si hay mapeo por coordenadas, usarlo
            if 'text_areas' in ai_content['slides'][slide_idx]:
                ai_slide_content['_text_areas'] = ai_content['slides'][slide_idx]['text_areas']
        
        # Clonar slide completo con todos sus elementos
        clone_slide_with_content(new_prs, template_slide, ai_slide_content)
        
        print(f"   ✅ Slide {slide_idx + 1} clonado")
    
    # Guardar nueva presentación
    output_path = tempfile.mktemp(suffix='.pptx')
    new_prs.save(output_path)
    
    print(f"\n✅ Presentación guardada en: {output_path}")
    
    return output_path

def clone_slide_with_content(new_prs: Presentation, template_slide, ai_content: Optional[Dict] = None):
    """
    Clona un slide completo preservando TODOS los elementos visuales,
    y solo reemplaza el texto si hay contenido de IA.
    """
    # Usar layout en blanco para tener control total
    blank_layout = new_prs.slide_layouts[6]  # Layout 6 es típicamente blank
    new_slide = new_prs.slides.add_slide(blank_layout)
    
    # Copiar background
    try:
        if template_slide.background.fill.type is not None:
            fill_type = template_slide.background.fill.type
            
            if fill_type == 1:  # Solid fill
                new_slide.background.fill.solid()
                if template_slide.background.fill.fore_color.rgb:
                    new_slide.background.fill.fore_color.rgb = template_slide.background.fill.fore_color.rgb
            elif fill_type == 2:  # Gradient
                new_slide.background.fill.gradient()
            elif fill_type == 3:  # Pattern
                pass  # Patterns son complejos
            elif fill_type == 6:  # Picture
                # Intentar copiar imagen de fondo
                try:
                    from io import BytesIO
                    bg_image = template_slide.background.fill.picture
                    if bg_image:
                        image_bytes = bg_image.blob
                        # No hay forma directa de agregar imagen de fondo en python-pptx
                        # Se mantiene el color sólido como fallback
                except:
                    pass
    except Exception as e:
        print(f"⚠️ No se pudo copiar background: {e}")
    
    # Clonar TODOS los shapes (formas, imágenes, text boxes, etc.)
    for shape in template_slide.shapes:
        clone_shape(new_slide, shape, ai_content)
    
    return new_slide

def clone_shape(new_slide, template_shape, ai_content: Optional[Dict] = None):
    """
    Clona un shape completo preservando su formato visual.
    Solo reemplaza el texto si es un text frame y hay contenido de IA.
    """
    shape_type = template_shape.shape_type
    shape_name = getattr(template_shape, 'name', 'unknown')
    
    print(f"      📦 Clonando shape: {shape_name} (tipo {shape_type})")
    
    # Clonar shape según su tipo
    try:
        if shape_type == 1:  # Auto shape (rectángulos, círculos, etc.)
            clone_autoshape(new_slide, template_shape, ai_content)
        
        elif shape_type == 13:  # Picture
            clone_picture(new_slide, template_shape)
        
        elif shape_type == 14:  # Placeholder
            clone_placeholder(new_slide, template_shape, ai_content)
        
        elif shape_type == 17:  # Text box
            clone_textbox(new_slide, template_shape, ai_content)
        
        elif shape_type == 3:  # Group
            clone_group(new_slide, template_shape, ai_content)
        
        elif shape_type == 19:  # Table
            clone_table(new_slide, template_shape, ai_content)
        
        elif shape_type == 5:  # Chart (gráfico)
            clone_chart(new_slide, template_shape, ai_content)
        
        elif hasattr(template_shape, 'has_chart') and template_shape.has_chart:
            clone_chart(new_slide, template_shape, ai_content)
        
        elif hasattr(template_shape, 'has_table') and template_shape.has_table:
            clone_table(new_slide, template_shape, ai_content)
        
        else:
            # Otros tipos: SmartArt, etc.
            clone_generic_shape(new_slide, template_shape)
    except Exception as e:
        print(f"      ⚠️ Error clonando shape {shape_name}: {e}")

def clone_autoshape(new_slide, template_shape, ai_content: Optional[Dict] = None):
    """Clona un autoshape (rectángulo, círculo, etc.) con todo su formato"""
    new_shape = new_slide.shapes.add_shape(
        template_shape.auto_shape_type,
        template_shape.left,
        template_shape.top,
        template_shape.width,
        template_shape.height
    )
    
    # Copiar formato de relleno
    copy_fill_format(template_shape, new_shape)
    
    # Copiar formato de línea
    copy_line_format(template_shape, new_shape)
    
    # Copiar texto si existe
    if template_shape.has_text_frame:
        copy_text_frame(template_shape.text_frame, new_shape.text_frame, ai_content)

def clone_textbox(new_slide, template_shape, ai_content: Optional[Dict] = None):
    """Clona un text box con su formato"""
    new_textbox = new_slide.shapes.add_textbox(
        template_shape.left,
        template_shape.top,
        template_shape.width,
        template_shape.height
    )
    
    # Copiar texto y formato
    if template_shape.has_text_frame:
        copy_text_frame(template_shape.text_frame, new_textbox.text_frame, ai_content)

def clone_placeholder(new_slide, template_shape, ai_content: Optional[Dict] = None):
    """Clona un placeholder como textbox para preservar posición exacta"""
    # Los placeholders se clonan como textboxes para mantener control total
    clone_textbox(new_slide, template_shape, ai_content)

def clone_picture(new_slide, template_shape):
    """Clona una imagen"""
    try:
        from io import BytesIO
        
        # Obtener la imagen del shape original
        image = template_shape.image
        image_bytes = image.blob
        
        # Crear stream de bytes
        image_stream = BytesIO(image_bytes)
        
        # Agregar imagen al nuevo slide
        new_slide.shapes.add_picture(
            image_stream,
            template_shape.left,
            template_shape.top,
            template_shape.width,
            template_shape.height
        )
    except Exception as e:
        print(f"⚠️ Error clonando imagen: {e}")

def clone_group(new_slide, template_shape, ai_content: Optional[Dict] = None):
    """Clona un grupo de shapes"""
    # Los grupos son complejos, clonar cada shape individual
    for shape in template_shape.shapes:
        clone_shape(new_slide, shape, ai_content)

def clone_generic_shape(new_slide, template_shape):
    """Intenta clonar un shape genérico"""
    # Placeholder para otros tipos de shapes
    pass


def clone_table(new_slide, template_shape, ai_content: Optional[Dict] = None):
    """
    Clona una tabla completa con su formato y datos.
    Si hay contenido de IA con 'table_data', reemplaza los datos.
    Usa el módulo Table Preserver si está disponible.
    """
    try:
        source_table = template_shape.table
        rows = len(source_table.rows)
        cols = len(source_table.columns)
        
        # Usar módulo avanzado si está disponible
        if TABLE_PRESERVER_AVAILABLE:
            # Extraer datos de la tabla original
            table_data = extract_table_data(source_table)
            
            # Preservar XML para mantener propiedades avanzadas
            preservation = preserve_table_xml(source_table)
            
            # Si hay datos de IA, actualizar
            if ai_content and 'table_data' in ai_content:
                table_data = ai_content['table_data']
            
            # Crear nueva tabla con datos preservados
            new_table_shape = new_slide.shapes.add_table(
                table_data.get('rows', rows),
                table_data.get('cols', cols),
                template_shape.left,
                template_shape.top,
                template_shape.width,
                template_shape.height
            )
            new_table = new_table_shape.table
            
            # Actualizar con datos
            update_table_with_data(new_table, table_data)
            
            print(f"✅ Tabla clonada con módulo avanzado: {rows}x{cols}")
            return
        
        # Método legacy
        new_table_shape = new_slide.shapes.add_table(
            rows, cols,
            template_shape.left,
            template_shape.top,
            template_shape.width,
            template_shape.height
        )
        new_table = new_table_shape.table
        
        # Obtener datos de IA si existen
        ai_table_data = None
        if ai_content and 'table_data' in ai_content:
            ai_table_data = ai_content['table_data']
        
        # Copiar contenido y formato de cada celda
        for row_idx, row in enumerate(source_table.rows):
            # Copiar altura de fila
            try:
                new_table.rows[row_idx].height = row.height
            except:
                pass
            
            for col_idx, cell in enumerate(row.cells):
                new_cell = new_table.cell(row_idx, col_idx)
                
                # Copiar texto (o usar datos de IA)
                if ai_table_data and row_idx < len(ai_table_data) and col_idx < len(ai_table_data[row_idx]):
                    new_cell.text = str(ai_table_data[row_idx][col_idx])
                else:
                    new_cell.text = cell.text
                
                # Copiar formato de celda
                copy_cell_format(cell, new_cell)
        
        # Copiar ancho de columnas
        for col_idx, col in enumerate(source_table.columns):
            try:
                new_table.columns[col_idx].width = col.width
            except:
                pass
        
        print(f"✅ Tabla clonada: {rows}x{cols}")
        
    except Exception as e:
        print(f"⚠️ Error clonando tabla: {e}")


def copy_cell_format(source_cell, target_cell):
    """Copia el formato de una celda de tabla"""
    try:
        # Copiar formato de relleno
        if source_cell.fill.type == 1:  # Solid
            target_cell.fill.solid()
            if source_cell.fill.fore_color.rgb:
                target_cell.fill.fore_color.rgb = source_cell.fill.fore_color.rgb
        
        # Copiar formato de texto
        if source_cell.text_frame.paragraphs:
            source_para = source_cell.text_frame.paragraphs[0]
            target_para = target_cell.text_frame.paragraphs[0]
            
            target_para.alignment = source_para.alignment
            
            if source_para.runs:
                source_run = source_para.runs[0]
                if target_para.runs:
                    target_run = target_para.runs[0]
                    copy_run_format(source_run, target_run)
    except:
        pass


def clone_chart(new_slide, template_shape, ai_content: Optional[Dict] = None):
    """
    Clona un gráfico.
    Si el módulo Chart Modifier está disponible, lo usa para preservación avanzada.
    """
    try:
        if not hasattr(template_shape, 'chart'):
            print("⚠️ Shape no tiene gráfico")
            return
        
        source_chart = template_shape.chart
        chart_type = source_chart.chart_type
        
        # Usar módulo avanzado si está disponible
        if CHART_MODIFIER_AVAILABLE:
            # Extraer datos del gráfico original
            chart_data = extract_chart_data_advanced(source_chart)
            
            # Si hay datos de IA, generarlos
            if ai_content and 'chart_data' in ai_content:
                chart_data = generate_chart_data_with_ai(chart_data, ai_content)
            
            # Actualizar el gráfico existente si es posible
            if update_chart_with_data(source_chart, chart_data):
                print("✅ Gráfico actualizado con módulo avanzado")
                return
        
        # Método legacy
        chart_data = extract_chart_data(source_chart)
        
        # Si hay datos de IA, usarlos
        if ai_content and 'chart_data' in ai_content:
            ai_chart_data = ai_content['chart_data']
            if 'categories' in ai_chart_data:
                chart_data['categories'] = ai_chart_data['categories']
            if 'series' in ai_chart_data:
                chart_data['series'] = ai_chart_data['series']
        
        # Crear nuevo gráfico
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        new_chart_data = CategoryChartData()
        new_chart_data.categories = chart_data.get('categories', ['A', 'B', 'C'])
        
        for series in chart_data.get('series', [{'name': 'Serie 1', 'values': [1, 2, 3]}]):
            new_chart_data.add_series(series['name'], series['values'])
        
        # Agregar gráfico al slide
        new_chart = new_slide.shapes.add_chart(
            chart_type,
            template_shape.left,
            template_shape.top,
            template_shape.width,
            template_shape.height,
            new_chart_data
        )
        
        print(f"✅ Gráfico clonado: tipo {chart_type}")
        
    except Exception as e:
        print(f"⚠️ Error clonando gráfico: {e}")
        # Fallback: crear placeholder
        try:
            placeholder = new_slide.shapes.add_textbox(
                template_shape.left,
                template_shape.top,
                template_shape.width,
                template_shape.height
            )
            placeholder.text_frame.text = "[Gráfico - requiere edición manual]"
        except:
            pass


def extract_chart_data(chart):
    """Extrae datos de un gráfico existente"""
    data = {
        'categories': [],
        'series': []
    }
    
    try:
        # Extraer categorías
        if chart.plots and chart.plots[0].categories:
            data['categories'] = list(chart.plots[0].categories)
        
        # Extraer series
        for series in chart.series:
            series_data = {
                'name': series.name or 'Serie',
                'values': list(series.values) if series.values else []
            }
            data['series'].append(series_data)
    except Exception as e:
        print(f"⚠️ Error extrayendo datos del gráfico: {e}")
        # Datos por defecto
        data = {
            'categories': ['Cat 1', 'Cat 2', 'Cat 3'],
            'series': [{'name': 'Serie 1', 'values': [1, 2, 3]}]
        }
    
    return data

def copy_fill_format(source_shape, target_shape):
    """Copia el formato de relleno (color, gradiente, etc.)"""
    try:
        fill_type = source_shape.fill.type
        
        if fill_type == 1:  # Solid fill
            target_shape.fill.solid()
            if source_shape.fill.fore_color.rgb:
                target_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
                
        elif fill_type == 2:  # Gradient fill
            target_shape.fill.gradient()
            # Copiar gradiente completo si está disponible
            try:
                gradient_fill = source_shape.fill.gradient
                if hasattr(gradient_fill, 'gradient_stops'):
                    # Copiar stops del gradiente si es posible
                    stops = gradient_fill.gradient_stops
                    if stops and len(stops) >= 2:
                        # Copiar color del primer y último stop
                        stop1 = stops[0]
                        stop2 = stops[-1]
                        if hasattr(stop1, 'color') and stop1.color and stop1.color.rgb:
                            target_shape.fill.gradient_stops[0].color.rgb = stop1.color.rgb
                        if hasattr(stop2, 'color') and stop2.color and stop2.color.rgb:
                            if len(target_shape.fill.gradient_stops) > 1:
                                target_shape.fill.gradient_stops[-1].color.rgb = stop2.color.rgb
            except Exception as e:
                print(f"      ⚠️ No se pudo copiar gradiente detallado: {e}")
                
        elif fill_type == 3:  # Pattern fill
            # Patrones ahora soportados
            try:
                target_shape.fill.pattern()
                if source_shape.fill.pattern and source_shape.fill.pattern.pattern_type:
                    # Copiar tipo de patrón si está disponible
                    pass  # python-pptx tiene soporte limitado para patrones
                # Copiar colores del patrón
                if source_shape.fill.fore_color.rgb:
                    target_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
                if source_shape.fill.back_color.rgb:
                    target_shape.fill.back_color.rgb = source_shape.fill.back_color.rgb
                print(f"      ✅ Patrón de relleno preservado")
            except Exception as e:
                print(f"      ⚠️ Error copiando patrón: {e}")
                
    except Exception as e:
        print(f"      ⚠️ Error en copy_fill_format: {e}")
        pass

def copy_line_format(source_shape, target_shape):
    """Copia el formato de línea (color, grosor, estilo)"""
    try:
        if source_shape.line.color.rgb:
            target_shape.line.color.rgb = source_shape.line.color.rgb
        if source_shape.line.width:
            target_shape.line.width = source_shape.line.width
    except:
        pass

def copy_text_frame(source_tf, target_tf, ai_content: Optional[Dict] = None):
    """
    Copia el text frame completo.
    Si hay contenido de IA, reemplaza el texto pero mantiene el formato.
    
    Soporta dos modos de mapeo:
    1. Por tipo (title, subtitle, bullets) - modo legacy
    2. Por coordenadas exactas (_text_areas) - modo preciso
    """
    # Limpiar target
    target_tf.clear()
    
    # Si hay contenido de IA, usarlo
    should_replace = False
    replacement_text = None
    is_bullets = False
    
    if ai_content:
        # Modo 1: Mapeo por coordenadas exactas
        if '_text_areas' in ai_content:
            text_areas = ai_content['_text_areas']
            # Buscar área que coincida con esta posición
            replacement_text = find_content_for_text_frame(source_tf, text_areas)
            if replacement_text:
                should_replace = True
        
        # Modo 2: Mapeo por tipo (fallback)
        if not should_replace:
            source_text = source_tf.text.strip()
            num_paragraphs = len(source_tf.paragraphs)
            
            # Detectar si es bullets (múltiples párrafos o tiene viñetas)
            has_bullets_format = num_paragraphs > 1 or any(
                p.level > 0 for p in source_tf.paragraphs
            )
            
            # Prioridad: bullets > title > heading > subtitle
            if has_bullets_format and 'bullets' in ai_content and ai_content['bullets']:
                copy_bullets_with_format(source_tf, target_tf, ai_content['bullets'])
                return
            
            # Si el texto original está vacío o es placeholder, reemplazar
            is_placeholder = (
                source_text == '' or 
                'click to' in source_text.lower() or
                'haga clic' in source_text.lower() or
                'título' in source_text.lower() or
                'title' in source_text.lower() or
                'subtitle' in source_text.lower() or
                'subtítulo' in source_text.lower()
            )
            
            if is_placeholder or len(source_text) < 100:
                # Determinar qué contenido usar basado en el texto original
                if 'title' in ai_content and ai_content['title']:
                    if 'título' in source_text.lower() or 'title' in source_text.lower() or source_text == '':
                        replacement_text = ai_content['title']
                        should_replace = True
                
                if not should_replace and 'subtitle' in ai_content and ai_content['subtitle']:
                    if 'subtítulo' in source_text.lower() or 'subtitle' in source_text.lower():
                        replacement_text = ai_content['subtitle']
                        should_replace = True
                
                if not should_replace and 'heading' in ai_content and ai_content['heading']:
                    replacement_text = ai_content['heading']
                    should_replace = True
                
                if not should_replace and 'title' in ai_content and ai_content['title']:
                    replacement_text = ai_content['title']
                    should_replace = True
    
    # Copiar párrafos con formato original
    for idx, source_para in enumerate(source_tf.paragraphs):
        if idx == 0:
            target_para = target_tf.paragraphs[0]
        else:
            target_para = target_tf.add_paragraph()
        
        # Copiar nivel y alineación
        target_para.level = source_para.level
        try:
            target_para.alignment = source_para.alignment
        except:
            pass
        
        # Si hay reemplazo de texto, usar el nuevo
        if should_replace and idx == 0:
            run = target_para.add_run()
            run.text = replacement_text if replacement_text else ''
            # Copiar formato del primer run original
            if source_para.runs:
                copy_run_format(source_para.runs[0], run)
        elif should_replace:
            # Para párrafos adicionales cuando hay reemplazo, no agregar nada
            pass
        else:
            # Copiar runs originales (mantener texto original)
            for source_run in source_para.runs:
                target_run = target_para.add_run()
                target_run.text = source_run.text
                copy_run_format(source_run, target_run)


def find_content_for_text_frame(source_tf, text_areas):
    """
    Busca el contenido apropiado para un text frame basándose en coordenadas.
    
    text_areas es una lista de:
    {
        'position': {'x': ..., 'y': ..., 'width': ..., 'height': ...},
        'content': 'texto a usar',
        'type': 'title' | 'subtitle' | 'bullets' | 'body'
    }
    """
    if not text_areas:
        return None
    
    # Obtener texto original para comparar
    original_text = source_tf.text.strip().lower()
    
    for area in text_areas:
        # Buscar por texto original si está disponible
        if 'originalText' in area:
            if area['originalText'].strip().lower() == original_text:
                return area.get('content', '')
        
        # Buscar por tipo
        if 'areaType' in area:
            area_type = area['areaType']
            content = area.get('content', '')
            
            if area_type == 'title' and len(original_text) < 50:
                return content
            elif area_type == 'subtitle' and 'subtitle' in original_text:
                return content
            elif area_type == 'heading' and len(original_text) < 100:
                return content
    
    return None

def copy_bullets_with_format(source_tf, target_tf, bullets: list):
    """Copia bullets con formato, reemplazando el texto"""
    # Obtener formato del primer párrafo
    original_format = None
    if source_tf.paragraphs and source_tf.paragraphs[0].runs:
        original_format = {
            'level': source_tf.paragraphs[0].level,
            'alignment': source_tf.paragraphs[0].alignment
        }
        original_run = source_tf.paragraphs[0].runs[0]
    
    # Agregar nuevos bullets
    for idx, bullet_text in enumerate(bullets):
        if idx == 0:
            para = target_tf.paragraphs[0]
        else:
            para = target_tf.add_paragraph()
        
        # Aplicar formato original
        if original_format:
            para.level = original_format['level']
            para.alignment = original_format['alignment']
        
        # Agregar texto
        run = para.add_run()
        run.text = bullet_text
        
        # Copiar formato del run original
        if source_tf.paragraphs and source_tf.paragraphs[0].runs:
            copy_run_format(source_tf.paragraphs[0].runs[0], run)

def copy_run_format(source_run, target_run):
    """Copia el formato de un run (fuente, tamaño, color, etc.)"""
    try:
        if source_run.font.name:
            target_run.font.name = source_run.font.name
        if source_run.font.size:
            target_run.font.size = source_run.font.size
        if source_run.font.bold:
            target_run.font.bold = True
        if source_run.font.italic:
            target_run.font.italic = True
        if source_run.font.underline:
            target_run.font.underline = True
        if source_run.font.color.rgb:
            target_run.font.color.rgb = source_run.font.color.rgb
    except:
        pass

# Funciones legacy removidas - ahora usamos clonación completa

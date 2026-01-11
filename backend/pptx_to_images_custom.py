"""
Renderizador personalizado de PPT a imágenes usando python-pptx + Pillow
100% gratis, sin marcas de agua, alta fidelidad
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import io
import base64
from typing import List

def convert_pptx_to_images_custom(pptx_path: str) -> List[str]:
    """
    Convierte cada slide del PPT a imagen base64 usando renderizado personalizado
    """
    try:
        prs = Presentation(pptx_path)
        base64_images = []
        
        # Obtener dimensiones del slide
        slide_width = int(prs.slide_width / 9525)  # EMUs a píxeles
        slide_height = int(prs.slide_height / 9525)
        
        # Escalar para mejor calidad
        scale = 2
        width = slide_width * scale
        height = slide_height * scale
        
        for slide_idx, slide in enumerate(prs.slides):
            # Crear imagen en blanco
            img = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(img)
            
            # Renderizar fondo
            render_background(slide, img, draw, width, height)
            
            # Renderizar cada shape
            for shape in slide.shapes:
                try:
                    render_shape(shape, img, draw, scale)
                except Exception as e:
                    print(f"⚠️ Error renderizando shape: {e}")
                    continue
            
            # Convertir a base64
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='PNG', quality=95)
            img_byte_arr.seek(0)
            
            img_str = base64.b64encode(img_byte_arr.getvalue()).decode()
            base64_images.append(f"data:image/png;base64,{img_str}")
            
            print(f"✅ Slide {slide_idx + 1} renderizado")
        
        return base64_images
        
    except Exception as e:
        print(f"❌ Error en renderizado custom: {e}")
        import traceback
        traceback.print_exc()
        return []

def render_background(slide, img, draw, width, height):
    """Renderiza el fondo del slide"""
    try:
        if slide.background.fill.type == 1:  # Solid fill
            color = get_rgb_color(slide.background.fill.fore_color)
            if color:
                img.paste(color, [0, 0, width, height])
    except:
        pass  # Fondo blanco por defecto

def render_shape(shape, img, draw, scale):
    """Renderiza un shape (texto, imagen, forma)"""
    
    # Convertir posición de EMUs a píxeles
    x = int(shape.left / 9525 * scale)
    y = int(shape.top / 9525 * scale)
    w = int(shape.width / 9525 * scale)
    h = int(shape.height / 9525 * scale)
    
    # Renderizar según tipo
    if shape.has_text_frame:
        render_text_shape(shape, draw, x, y, w, h, scale)
    
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        render_image_shape(shape, img, x, y, w, h)
    
    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        render_auto_shape(shape, draw, x, y, w, h)

def render_text_shape(shape, draw, x, y, w, h, scale):
    """Renderiza un shape con texto"""
    text_frame = shape.text_frame
    
    # Renderizar fondo del shape si tiene
    try:
        if shape.fill.type == 1:  # Solid fill
            color = get_rgb_color(shape.fill.fore_color)
            if color:
                draw.rectangle([x, y, x + w, y + h], fill=color)
    except:
        pass
    
    # Renderizar cada párrafo
    current_y = y + 10
    
    for paragraph in text_frame.paragraphs:
        if not paragraph.text.strip():
            continue
            
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            
            # Obtener formato del texto
            font_name = run.font.name or "Arial"
            font_size = int((run.font.size.pt if run.font.size else 18) * scale)
            
            try:
                font = ImageFont.truetype(f"C:/Windows/Fonts/{font_name}.ttf", font_size)
            except:
                try:
                    font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", font_size)
                except:
                    font = ImageFont.load_default()
            
            # Color del texto
            text_color = get_rgb_color(run.font.color) if run.font.color else (0, 0, 0)
            
            # Dibujar texto
            text = run.text
            
            # Calcular posición según alineación
            text_x = x + 20
            if paragraph.alignment == 2:  # Center
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_x = x + (w - text_width) // 2
            elif paragraph.alignment == 3:  # Right
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_x = x + w - text_width - 20
            
            draw.text((text_x, current_y), text, fill=text_color, font=font)
            
            # Avanzar línea
            bbox = draw.textbbox((0, 0), text, font=font)
            current_y += (bbox[3] - bbox[1]) + 10

def render_image_shape(shape, img, x, y, w, h):
    """Renderiza una imagen"""
    try:
        image_stream = shape.image.blob
        shape_img = Image.open(io.BytesIO(image_stream))
        shape_img = shape_img.resize((w, h), Image.Resampling.LANCZOS)
        img.paste(shape_img, (x, y))
    except Exception as e:
        print(f"⚠️ Error renderizando imagen: {e}")

def render_auto_shape(shape, draw, x, y, w, h):
    """Renderiza formas automáticas (rectángulos, círculos, etc)"""
    try:
        # Color de relleno
        fill_color = None
        if shape.fill.type == 1:  # Solid fill
            fill_color = get_rgb_color(shape.fill.fore_color)
        
        # Color de línea
        line_color = None
        if shape.line.fill.type == 1:
            line_color = get_rgb_color(shape.line.fill.fore_color)
        
        # Dibujar forma
        draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline=line_color, width=2)
        
    except Exception as e:
        print(f"⚠️ Error renderizando forma: {e}")

def get_rgb_color(color_obj):
    """Obtiene color RGB de un objeto de color de python-pptx"""
    try:
        if hasattr(color_obj, 'rgb') and color_obj.rgb:
            rgb = color_obj.rgb
            return (rgb[0], rgb[1], rgb[2])
    except:
        pass
    return None

"""
Script de diagnóstico para ver el XML real del fondo de un slide
"""
from pptx import Presentation
from lxml import etree
import sys

def debug_slide_background(pptx_path, slide_number):
    """Muestra el XML del fondo de un slide específico"""
    prs = Presentation(pptx_path)
    
    if slide_number > len(prs.slides):
        print(f"❌ El PPTX solo tiene {len(prs.slides)} slides")
        return
    
    slide = prs.slides[slide_number - 1]
    
    print(f"\n{'='*80}")
    print(f"DIAGNÓSTICO DE FONDO - SLIDE {slide_number}")
    print(f"{'='*80}\n")
    
    # Obtener XML del slide
    slide_part = slide.part
    slide_xml = slide_part.blob
    root = etree.fromstring(slide_xml)
    
    # Namespaces
    namespaces = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    
    print("📄 XML DEL SLIDE (fragmento de fondo):")
    print("-" * 80)
    
    # Buscar elemento de fondo
    bg_element = root.find('.//p:cSld/p:bg', namespaces)
    
    if bg_element is not None:
        print("✅ Elemento p:bg encontrado en el slide:")
        print(etree.tostring(bg_element, pretty_print=True, encoding='unicode'))
    else:
        print("❌ No se encontró elemento p:bg en el slide")
        print("\n🔍 Buscando en el layout...")
        
        # Buscar en layout
        layout_part = slide.slide_layout.part
        layout_xml = layout_part.blob
        layout_root = etree.fromstring(layout_xml)
        
        bg_element = layout_root.find('.//p:cSld/p:bg', namespaces)
        
        if bg_element is not None:
            print("✅ Elemento p:bg encontrado en el layout:")
            print(etree.tostring(bg_element, pretty_print=True, encoding='unicode'))
        else:
            print("❌ No se encontró elemento p:bg en el layout")
            print("\n🔍 Buscando en el master...")
            
            # Buscar en master
            try:
                master_part = slide.slide_layout.slide_master.part
                master_xml = master_part.blob
                master_root = etree.fromstring(master_xml)
                
                bg_element = master_root.find('.//p:cSld/p:bg', namespaces)
                
                if bg_element is not None:
                    print("✅ Elemento p:bg encontrado en el master:")
                    print(etree.tostring(bg_element, pretty_print=True, encoding='unicode'))
                else:
                    print("❌ No se encontró elemento p:bg en el master")
            except Exception as e:
                print(f"⚠️ Error accediendo al master: {e}")
    
    print("\n" + "="*80)
    print("📊 INFORMACIÓN DEL TEMA")
    print("="*80 + "\n")
    
    # Intentar obtener información del tema
    try:
        master = slide.slide_layout.slide_master
        theme_part = master.part.part_related_by('http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
        theme_xml = theme_part.blob
        theme_root = etree.fromstring(theme_xml)
        
        print("🎨 Tema encontrado:")
        
        # Buscar color scheme
        color_scheme = theme_root.find('.//a:clrScheme', namespaces)
        if color_scheme is not None:
            print("\n📋 Esquema de colores:")
            for color_elem in color_scheme:
                color_name = color_elem.tag.split('}')[-1]
                
                # Buscar el valor del color
                srgb = color_elem.find('.//a:srgbClr', namespaces)
                sys_clr = color_elem.find('.//a:sysClr', namespaces)
                
                if srgb is not None:
                    color_val = srgb.get('val')
                    print(f"   {color_name}: #{color_val}")
                elif sys_clr is not None:
                    color_val = sys_clr.get('lastClr')
                    print(f"   {color_name}: #{color_val} (system)")
                else:
                    print(f"   {color_name}: (complejo)")
    except Exception as e:
        print(f"⚠️ No se pudo acceder al tema: {e}")
    
    print("\n" + "="*80 + "\n")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python debug_slide_background.py <archivo.pptx> [numero_slide]")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2]) if len(sys.argv) > 2 else 2
    
    debug_slide_background(pptx_path, slide_num)
